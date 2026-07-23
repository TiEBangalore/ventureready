"""
VentureReady demo app server.
Serves the prototype (index.html) and provides two LIVE endpoints:
  - POST /api/verify-member  -> checks an email against TiE Bangalore's Zoho membership data
  - POST /api/diagnostic     -> runs the AI positioning read (uses Claude if a key is set, else a canned result)
Run:  python3 server.py    then open http://localhost:8000
"""
import json, time, os, io, base64, sqlite3, urllib.parse, urllib.request, urllib.error
import hashlib, hmac, secrets, re
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer

HERE = os.path.dirname(os.path.abspath(__file__))

# ---- Local database (SQLite). Proves real persistence: data survives refresh AND restart. ----
DB_PATH = os.path.join(HERE, "data.db")
def _db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn
def db_init():
    conn = _db()
    conn.execute("""CREATE TABLE IF NOT EXISTS support(
        id INTEGER PRIMARY KEY AUTOINCREMENT, ref TEXT, from_name TEXT, role TEXT,
        category TEXT, subject TEXT, message TEXT, status TEXT DEFAULT 'New',
        received TEXT, created_at TEXT)""")
    # Each founder is one row. email is unique so the same person maps to one record.
    conn.execute("""CREATE TABLE IF NOT EXISTS founder(
        id INTEGER PRIMARY KEY AUTOINCREMENT, name TEXT, email TEXT UNIQUE, company TEXT,
        role TEXT, phone TEXT, city TEXT, stage TEXT, sector TEXT, linkedin TEXT,
        created_at TEXT)""")
    # Migration: older databases don't have password columns. SQLite has no
    # "ADD COLUMN IF NOT EXISTS", so we check the existing columns first.
    _cols = [r["name"] for r in conn.execute("PRAGMA table_info(founder)").fetchall()]
    if "password_hash" not in _cols:
        conn.execute("ALTER TABLE founder ADD COLUMN password_hash TEXT")
    if "password_salt" not in _cols:
        conn.execute("ALTER TABLE founder ADD COLUMN password_salt TEXT")
    # Membership snapshot columns. These CACHE Zoho's answer, refreshed at each
    # sign-in (Zoho stays the source of truth). membership_stale=1 means the last
    # refresh couldn't reach Zoho, so these values are the last CONFIRMED status.
    for _mcol, _mtype in (("is_member", "INTEGER"), ("membership_status", "TEXT"),
                          ("membership_category", "TEXT"), ("membership_name", "TEXT"),
                          ("membership_checked_at", "TEXT"), ("membership_stale", "INTEGER")):
        if _mcol not in _cols:
            conn.execute("ALTER TABLE founder ADD COLUMN %s %s" % (_mcol, _mtype))
    # google_sub = Google's stable, unique ID for the account (the "sub" claim).
    # We match on this first so a founder who later changes their Gmail address
    # still maps to the same record. Nullable: password-only accounts don't have it.
    if "google_sub" not in _cols:
        conn.execute("ALTER TABLE founder ADD COLUMN google_sub TEXT")
    # Each uploaded deck: the file lives on disk (stored_path); the DB keeps a pointer + owner.
    conn.execute("""CREATE TABLE IF NOT EXISTS deck(
        id INTEGER PRIMARY KEY AUTOINCREMENT, founder_id INTEGER, filename TEXT,
        stored_path TEXT, size INTEGER, uploaded_at TEXT)""")
    # Each AI read is saved against the founder + the deck it read.
    conn.execute("""CREATE TABLE IF NOT EXISTS diagnostic(
        id INTEGER PRIMARY KEY AUTOINCREMENT, founder_id INTEGER, deck_id INTEGER,
        live INTEGER, summary TEXT, findings_json TEXT, created_at TEXT)""")
    # Expert review rounds. round 1 = the paid review; round 2 = the ONE free
    # re-review; round 3+ = paid again. verdict: 'not_yet' | 'awarded'.
    conn.execute("""CREATE TABLE IF NOT EXISTS review(
        id INTEGER PRIMARY KEY AUTOINCREMENT, founder_id INTEGER, round INTEGER,
        verdict TEXT, gaps_json TEXT, reviewer TEXT, note TEXT, created_at TEXT)""")
    # Data room: one row per checklist item a founder has supplied a document for.
    # item_key ties the file back to the diligence checklist in the front-end.
    conn.execute("""CREATE TABLE IF NOT EXISTS dataroom_doc(
        id INTEGER PRIMARY KEY AUTOINCREMENT, founder_id INTEGER, item_key TEXT,
        filename TEXT, stored_path TEXT, size INTEGER, uploaded_at TEXT)""")
    # Who opened which data-room document — the "who's actually interested" signal.
    conn.execute("""CREATE TABLE IF NOT EXISTS dataroom_view(
        id INTEGER PRIMARY KEY AUTOINCREMENT, founder_id INTEGER, doc_id INTEGER,
        item_key TEXT, viewer TEXT, viewed_at TEXT)""")
    # Admin allow-list: who may enter the Team/Admin portal. A super-admin edits
    # this from inside the portal, so this table is the LIVE source of truth. On
    # first run it is seeded from the super-admins + ADMIN_EMAILS (.env).
    conn.execute("""CREATE TABLE IF NOT EXISTS admin_allowlist(
        email TEXT PRIMARY KEY, added_by TEXT, added_at TEXT)""")
    _admin_seed = set(ADMIN_SUPER_EMAILS)
    for _e in ENV.get("ADMIN_EMAILS", "").split(","):
        _e = _e.strip().lower()
        if _e:
            _admin_seed.add(_e)
    for _e in _admin_seed:
        conn.execute("INSERT OR IGNORE INTO admin_allowlist(email, added_by, added_at) "
                     "VALUES(?,?,datetime('now'))", (_e, "seed"))
    # Notifications: one row per important event, recorded even when email is off.
    conn.execute("""CREATE TABLE IF NOT EXISTS notification(
        id INTEGER PRIMARY KEY AUTOINCREMENT, event TEXT, recipient TEXT,
        subject TEXT, body TEXT, urgency TEXT, created_at TEXT)""")
    # Email outbox: one row per recipient. status: 'pending' | 'sent' | 'failed'.
    # In record-only mode (no SMTP configured) rows simply stay 'pending'.
    conn.execute("""CREATE TABLE IF NOT EXISTS email_outbox(
        id INTEGER PRIMARY KEY AUTOINCREMENT, notification_id INTEGER, to_email TEXT,
        subject TEXT, body TEXT, status TEXT DEFAULT 'pending', error TEXT,
        created_at TEXT, sent_at TEXT)""")
    # Payments: a real row every time a founder/investor pays for something
    # (membership, expert review, deal-flow access). No card data is stored —
    # this is the prototype record that drives the receipt notification.
    conn.execute("""CREATE TABLE IF NOT EXISTS payment(
        id INTEGER PRIMARY KEY AUTOINCREMENT, payer_email TEXT, payer_name TEXT,
        founder_id INTEGER, item TEXT, amount TEXT, period TEXT,
        status TEXT DEFAULT 'paid', created_at TEXT)""")
    # Review assignment: which reviewer is responsible for a founder's review,
    # who assigned it, and the SLA target. One row per assignment/reassignment.
    conn.execute("""CREATE TABLE IF NOT EXISTS review_assignment(
        id INTEGER PRIMARY KEY AUTOINCREMENT, founder_id INTEGER, founder_name TEXT,
        reviewer_email TEXT, reviewer_name TEXT, assigned_by TEXT,
        status TEXT DEFAULT 'assigned', sla_due TEXT, created_at TEXT)""")
    # Investors: real screened-access records. A primary applicant is a row with
    # parent_investor_id = NULL; each firm team-mate they invite is a row that
    # points back to the primary via parent_investor_id. status walks:
    # 'pending' -> 'approved' | 'rejected'; invited team-mates start 'invited'.
    conn.execute("""CREATE TABLE IF NOT EXISTS investor(
        id INTEGER PRIMARY KEY AUTOINCREMENT, name TEXT, email TEXT UNIQUE, firm TEXT,
        tie_status TEXT, cheque_size TEXT, focus TEXT, track_record TEXT,
        recent_investments TEXT, status TEXT DEFAULT 'pending', parent_investor_id INTEGER,
        decided_by TEXT, decided_at TEXT, created_at TEXT)""")
    # Meeting requests: an investor asks TiE to introduce them to a founder.
    # TiE facilitates every intro (no direct messaging). status walks:
    # 'requested' -> 'intro_sent' -> 'held' | 'declined'.
    conn.execute("""CREATE TABLE IF NOT EXISTS meeting(
        id INTEGER PRIMARY KEY AUTOINCREMENT, founder_id INTEGER, founder_name TEXT,
        startup TEXT, investor_id INTEGER, investor_name TEXT, investor_email TEXT,
        status TEXT DEFAULT 'requested', requested_at TEXT, intro_sent_at TEXT)""")
    # Reviewers (TiE Charter Members). Same auth shape as founders: password OR
    # Google. TiE controls who reviews by assigning decks to a reviewer's email.
    conn.execute("""CREATE TABLE IF NOT EXISTS reviewer(
        id INTEGER PRIMARY KEY AUTOINCREMENT, name TEXT, email TEXT UNIQUE,
        password_hash TEXT, password_salt TEXT, google_sub TEXT, created_at TEXT)""")
    # Review lifecycle: an AWARD is only a recommendation until an admin
    # countersigns it. Migrate older review rows (which had no status) to
    # 'issued' so nothing already decided is re-opened.
    _rcols = [r["name"] for r in conn.execute("PRAGMA table_info(review)").fetchall()]
    if "status" not in _rcols:
        conn.execute("ALTER TABLE review ADD COLUMN status TEXT DEFAULT 'issued'")
    if "countersigned_by" not in _rcols:
        conn.execute("ALTER TABLE review ADD COLUMN countersigned_by TEXT")
    if "countersigned_at" not in _rcols:
        conn.execute("ALTER TABLE review ADD COLUMN countersigned_at TEXT")
    # Deals: one founder<->investor relationship tracked from introduction to
    # outcome. This is the substrate for TiE's impact reporting. stage walks:
    # 'introduced' -> 'met' -> 'diligence' -> 'term_sheet' -> 'closed' | 'passed'.
    # amount_inr is optional (a deal can be 'closed' with amount undisclosed).
    # founder_confirmed / investor_confirmed record TWO-SIDED attestation of the
    # current stage; tie_verified is TiE's own spot-check. *_consent_public gate
    # whether the deal may EVER feed public/anonymised impact figures.
    conn.execute("""CREATE TABLE IF NOT EXISTS deal(
        id INTEGER PRIMARY KEY AUTOINCREMENT, founder_id INTEGER, founder_name TEXT,
        startup TEXT, investor_id INTEGER, investor_name TEXT, investor_email TEXT, firm TEXT,
        stage TEXT DEFAULT 'introduced', amount_inr INTEGER, amount_disclosed INTEGER DEFAULT 0,
        founder_consent_public INTEGER DEFAULT 0, investor_consent_public INTEGER DEFAULT 0,
        founder_confirmed INTEGER DEFAULT 0, investor_confirmed INTEGER DEFAULT 0,
        tie_verified INTEGER DEFAULT 0, source TEXT, meeting_id INTEGER, note TEXT,
        created_at TEXT, updated_at TEXT, closed_at TEXT)""")
    conn.commit()
    if conn.execute("SELECT COUNT(*) c FROM support").fetchone()["c"] == 0:
        seeds = [
            ("SUP-204", "Meera Suresh", "Founder", "TiE membership verification",
             "My member email isn't recognised at the gate", "", "New", "3h ago"),
            ("SUP-203", "Narendra Bhandari", "Investor", "Login or access problem",
             "Colleague can't accept the firm invite", "", "In progress", "Yesterday"),
            ("SUP-202", "Rajan Kumar", "Reviewer", "My review status",
             "Which deck is next in my queue?", "", "Resolved", "2 days ago"),
        ]
        conn.executemany("INSERT INTO support(ref,from_name,role,category,subject,message,status,received,created_at) "
                         "VALUES(?,?,?,?,?,?,?,?,datetime('now'))", seeds)
        conn.commit()
    conn.close()
def support_add(d):
    conn = _db()
    cur = conn.execute("INSERT INTO support(ref,from_name,role,category,subject,message,status,received,created_at) "
                       "VALUES('',?,?,?,?,?,'New','Just now',datetime('now'))",
                       (d.get("from_name", ""), d.get("role", ""), d.get("category", ""),
                        d.get("subject", ""), d.get("message", "")))
    rid = cur.lastrowid
    ref = "SUP-%d" % (204 + rid)
    conn.execute("UPDATE support SET ref=? WHERE id=?", (ref, rid))
    conn.commit()
    row = dict(conn.execute("SELECT * FROM support WHERE id=?", (rid,)).fetchone())
    conn.close()
    notify("support_request", ["admins"],
           "New support request: " + (d.get("subject") or "(no subject)"),
           (d.get("from_name") or "Someone") + " (" + (d.get("role") or "user") +
           ") submitted a support request. " + (d.get("message") or ""), "normal")
    return row
def support_list():
    conn = _db()
    rows = [dict(r) for r in conn.execute("SELECT * FROM support ORDER BY id DESC").fetchall()]
    conn.close()
    return {"support": rows}

# ---- Founders (real per-user records) ----
DECKS_DIR = os.path.join(HERE, "decks")

# --- Password handling (PROTOTYPE-GRADE, not production security) ---
# Passwords are never stored or logged in plaintext. Each founder gets a random
# per-user salt; we store only the PBKDF2-SHA256 hash. A real deployment should
# replace this with a vetted auth service and add rate-limiting / email verify.
def _hash_pw(password, salt=None):
    if salt is None:
        salt = secrets.token_hex(16)
    dk = hashlib.pbkdf2_hmac("sha256", (password or "").encode("utf-8"),
                             salt.encode("utf-8"), 100000)
    return dk.hex(), salt

def _public_founder(row):
    """Strip secret columns so a founder's password hash never leaves the server."""
    if row is None:
        return None
    out = dict(row)
    out.pop("password_hash", None)
    out.pop("password_salt", None)
    return out

def founder_signup(d):
    """Register a founder with a password. Errors if the email already has one."""
    conn = _db()
    email = (d.get("email") or "").strip().lower()
    password = d.get("password") or ""
    if not email or not password:
        conn.close()
        return {"error": "Email and password are required."}
    existing = conn.execute("SELECT * FROM founder WHERE email=?", (email,)).fetchone()
    if existing and existing["password_hash"]:
        conn.close()
        return {"error": "That email is already registered. Please log in instead."}
    pw_hash, pw_salt = _hash_pw(password)
    fields = ("name", "company", "role", "phone", "city", "stage", "sector", "linkedin")
    vals = [d.get(f, "") for f in fields]
    if existing:
        conn.execute("UPDATE founder SET name=?,company=?,role=?,phone=?,city=?,stage=?,sector=?,linkedin=?,"
                     "password_hash=?,password_salt=? WHERE id=?",
                     tuple(vals) + (pw_hash, pw_salt, existing["id"]))
        fid = existing["id"]
    else:
        cur = conn.execute("INSERT INTO founder(name,email,company,role,phone,city,stage,sector,linkedin,"
                           "password_hash,password_salt,created_at) VALUES(?,?,?,?,?,?,?,?,?,?,?,datetime('now'))",
                           (d.get("name", ""), email) + tuple(vals[1:]) + (pw_hash, pw_salt))
        fid = cur.lastrowid
    conn.commit()
    row = conn.execute("SELECT * FROM founder WHERE id=?", (fid,)).fetchone()
    conn.close()
    # Check TiE membership with Zoho as part of signing up, so the new profile
    # already shows the correct member status.
    refreshed = refresh_membership(fid)
    _pub = _public_founder(refreshed or row) or {}
    notify("founder_welcome", [email],
           "Welcome to VentureReady",
           "Thanks for creating your VentureReady account, " + (_pub.get("name") or "founder") +
           ". You can now run the free positioning diagnostic and submit your deck for expert review.",
           "normal")
    return {"founder": _pub}

def founder_login(email, password):
    """Verify credentials in constant time. Returns the public profile or an error."""
    conn = _db()
    email = (email or "").strip().lower()
    row = conn.execute("SELECT * FROM founder WHERE email=?", (email,)).fetchone()
    conn.close()
    if not row or not row["password_hash"] or not row["password_salt"]:
        return {"error": "No account found for that email, or it has no password set."}
    calc, _ = _hash_pw(password, row["password_salt"])
    if not hmac.compare_digest(calc, row["password_hash"]):
        return {"error": "Incorrect email or password."}
    # Re-check TiE membership with Zoho on every sign-in and cache the result.
    refreshed = refresh_membership(row["id"])
    return {"founder": _public_founder(refreshed or row)}

def founder_upsert(d):
    """Create or update a founder by email, then return the full row (with its id)."""
    conn = _db()
    email = (d.get("email") or "").strip().lower()
    fields = ("name", "company", "role", "phone", "city", "stage", "sector", "linkedin")
    vals = [d.get(f, "") for f in fields]
    existing = None
    if email:
        existing = conn.execute("SELECT * FROM founder WHERE email=?", (email,)).fetchone()
    if existing:
        conn.execute("UPDATE founder SET name=?,company=?,role=?,phone=?,city=?,stage=?,sector=?,linkedin=? WHERE id=?",
                     tuple(vals) + (existing["id"],))
        fid = existing["id"]
    else:
        cur = conn.execute("INSERT INTO founder(name,email,company,role,phone,city,stage,sector,linkedin,created_at) "
                           "VALUES(?,?,?,?,?,?,?,?,?,datetime('now'))",
                           (d.get("name", ""), email) + tuple(vals[1:]))
        fid = cur.lastrowid
    conn.commit()
    row = _public_founder(conn.execute("SELECT * FROM founder WHERE id=?", (fid,)).fetchone())
    conn.close()
    return row

def founder_get(fid):
    conn = _db()
    r = conn.execute("SELECT * FROM founder WHERE id=?", (fid,)).fetchone()
    out = _public_founder(r) if r else None
    if out:
        dg = conn.execute("SELECT d.*, k.filename FROM diagnostic d LEFT JOIN deck k ON k.id=d.deck_id "
                          "WHERE d.founder_id=? ORDER BY d.id DESC LIMIT 1", (fid,)).fetchone()
        if dg:
            out["latest_diagnostic"] = {
                "live": bool(dg["live"]), "summary": dg["summary"],
                "findings": json.loads(dg["findings_json"] or "[]"), "filename": dg["filename"] or "",
            }
    conn.close()
    return out

def deck_add(founder_id, filename, raw):
    """Save the uploaded deck file to disk and record a row pointing to it."""
    if not os.path.isdir(DECKS_DIR):
        os.makedirs(DECKS_DIR)
    safe = "".join(c for c in (filename or "deck.pdf") if c.isalnum() or c in "._- ")
    stored = "f%s_%d_%s" % (founder_id or 0, int(time.time()), safe)
    path = os.path.join(DECKS_DIR, stored)
    with open(path, "wb") as fh:
        fh.write(raw)
    conn = _db()
    cur = conn.execute("INSERT INTO deck(founder_id,filename,stored_path,size,uploaded_at) "
                       "VALUES(?,?,?,?,datetime('now'))", (founder_id, filename, stored, len(raw)))
    did = cur.lastrowid
    conn.commit()
    conn.close()
    fb = _founder_brief(founder_id)
    notify("deck_submitted", ["admins"],
           "New deck submitted for review",
           (fb["name"] or "A founder") +
           ((" (" + fb["company"] + ")") if fb["company"] else "") +
           " uploaded a pitch deck: " + (filename or "deck") + ".", "normal")
    return did

# ---- Expert review rounds & the one-free-re-review rule ----
FREE_REREVIEW_ROUND = 2      # round 1 = paid review, round 2 = the single free re-review
REREVIEW_FEE = "₹3,000 + GST"

def review_add(founder_id, verdict, gaps, reviewer, note):
    """Record a reviewer's verdict as the next round. An AWARD is only a
    RECOMMENDATION ('recommended') until an admin countersigns it — the founder
    is not told and the mark is not live yet. A 'not yet' is coaching, so it is
    'issued' straight away and the founder sees the feedback."""
    status = "recommended" if verdict == "awarded" else "issued"
    conn = _db()
    last = conn.execute("SELECT MAX(round) m FROM review WHERE founder_id=?", (founder_id,)).fetchone()
    rnd = (last["m"] or 0) + 1
    conn.execute("INSERT INTO review(founder_id,round,verdict,gaps_json,reviewer,note,status,created_at) "
                 "VALUES(?,?,?,?,?,?,?,datetime('now'))",
                 (founder_id, rnd, verdict, json.dumps(gaps or []), reviewer or "TiE Reviewer", note or "", status))
    conn.commit()
    conn.close()
    fb = _founder_brief(founder_id)
    if verdict == "awarded":
        # Founder is NOT told yet — this is awaiting TiE's countersign.
        notify("verdict_pending_countersign", ["admins"],
               "Verdict to countersign: " + (fb["name"] or ("founder #" + str(founder_id))),
               (reviewer or "A reviewer") + " recommends awarding the VentureReady mark to " +
               (fb["name"] or ("founder #" + str(founder_id))) +
               ". Countersign in the Team Portal to issue it.", "high")
    else:
        if fb["email"]:
            notify("review_feedback", [fb["email"]],
                   "Your VentureReady review feedback is ready",
                   "Your expert review is complete, with feedback to act on before the next round. "
                   "Sign in to VentureReady to see the details.", "normal")
        notify("verdict_recorded", ["admins"],
               "A reviewer verdict was recorded",
               (reviewer or "A reviewer") + " recorded a 'not yet' verdict for " +
               (fb["name"] or ("founder #" + str(founder_id))) + ".", "normal")
    return review_state(founder_id)

def _founder_visible_verdict(row):
    """What the FOUNDER should see for a review row, honouring the countersign
    gate: a recommended award reads as 'in_review' (not the mark yet); an
    overridden award reads as 'not_yet' (withheld — back to the feedback path)."""
    if row is None:
        return None
    v, st = row.get("verdict"), (row.get("status") or "issued")
    if v == "awarded" and st == "recommended":
        return "in_review"
    if v == "awarded" and st == "overridden":
        return "not_yet"
    return v

def review_state(founder_id):
    """Everything the founder's screens need: history, latest verdict, and whether
    the next re-review is the free one or has to be paid for. Awards that haven't
    been countersigned are shown as 'in_review', never as the live mark."""
    conn = _db()
    rows = [dict(r) for r in conn.execute(
        "SELECT round,verdict,gaps_json,reviewer,note,status,created_at FROM review "
        "WHERE founder_id=? ORDER BY round", (founder_id,)).fetchall()]
    conn.close()
    for r in rows:
        r["gaps"] = json.loads(r.pop("gaps_json") or "[]")
        # Present the founder-visible verdict; keep the raw one for admins/debug.
        r["raw_verdict"] = r["verdict"]
        r["verdict"] = _founder_visible_verdict(r)
    latest = rows[-1] if rows else None
    rounds_done = len(rows)
    next_round = rounds_done + 1
    return {
        "rounds": rows,
        "latest": latest,
        "rounds_done": rounds_done,
        "next_round": next_round,
        # The single free re-review is round 2 and only if they haven't earned the mark.
        "free_rereview_available": (next_round == FREE_REREVIEW_ROUND
                                    and bool(latest) and latest["verdict"] == "not_yet"),
        "free_rereview_used": rounds_done >= FREE_REREVIEW_ROUND,
        "rereview_fee": REREVIEW_FEE,
    }

# ---- Reviewers (accounts, queue, gated verdict submission) -----------------
def _public_reviewer(row):
    if row is None:
        return None
    out = dict(row)
    out.pop("password_hash", None)
    out.pop("password_salt", None)
    return out

def reviewer_signup(d):
    email = (d.get("email") or "").strip().lower()
    password = d.get("password") or ""
    if not email or not password:
        return {"error": "Email and password are required."}
    conn = _db()
    existing = conn.execute("SELECT * FROM reviewer WHERE email=?", (email,)).fetchone()
    if existing and existing["password_hash"]:
        conn.close()
        return {"error": "That email is already registered. Please log in instead."}
    pw_hash, pw_salt = _hash_pw(password)
    if existing:
        conn.execute("UPDATE reviewer SET name=?, password_hash=?, password_salt=? WHERE id=?",
                     (d.get("name", ""), pw_hash, pw_salt, existing["id"]))
        rid = existing["id"]
    else:
        cur = conn.execute("INSERT INTO reviewer(name,email,password_hash,password_salt,created_at) "
                           "VALUES(?,?,?,?,datetime('now'))",
                           (d.get("name", ""), email, pw_hash, pw_salt))
        rid = cur.lastrowid
    conn.commit()
    row = conn.execute("SELECT * FROM reviewer WHERE id=?", (rid,)).fetchone()
    conn.close()
    return {"reviewer": _public_reviewer(row)}

def reviewer_login(email, password):
    conn = _db()
    email = (email or "").strip().lower()
    row = conn.execute("SELECT * FROM reviewer WHERE email=?", (email,)).fetchone()
    conn.close()
    if not row or not row["password_hash"] or not row["password_salt"]:
        return {"error": "No reviewer account found for that email, or it has no password set."}
    calc, _ = _hash_pw(password, row["password_salt"])
    if not hmac.compare_digest(calc, row["password_hash"]):
        return {"error": "Incorrect email or password."}
    return {"reviewer": _public_reviewer(row)}

def reviewer_google_login(id_token):
    v = verify_google_token(id_token)
    if v.get("error"):
        return {"error": v["error"]}
    sub, email, name = v["sub"], (v["email"] or "").strip().lower(), v["name"]
    conn = _db()
    row = conn.execute("SELECT * FROM reviewer WHERE google_sub=?", (sub,)).fetchone()
    if not row and email:
        row = conn.execute("SELECT * FROM reviewer WHERE email=?", (email,)).fetchone()
    if row:
        conn.execute("UPDATE reviewer SET google_sub=?, name=COALESCE(NULLIF(name,''), ?) WHERE id=?",
                     (sub, name or "", row["id"]))
        rid = row["id"]
    else:
        cur = conn.execute("INSERT INTO reviewer(name,email,google_sub,created_at) "
                           "VALUES(?,?,?,datetime('now'))", (name or "TiE Reviewer", email, sub))
        rid = cur.lastrowid
    conn.commit()
    row = conn.execute("SELECT * FROM reviewer WHERE id=?", (rid,)).fetchone()
    conn.close()
    return {"reviewer": _public_reviewer(row)}

def reviewer_queue(email):
    """A reviewer's own assigned decks: latest assignment per founder, joined with
    the founder brief and the current review status for that founder."""
    email = (email or "").strip().lower()
    if not email:
        return []
    conn = _db()
    rows = conn.execute(
        "SELECT a.founder_id, a.created_at AS assigned_at, a.sla_due, a.status AS assign_status "
        "FROM review_assignment a "
        "WHERE a.reviewer_email=? AND a.id IN "
        "  (SELECT MAX(id) FROM review_assignment WHERE reviewer_email=? GROUP BY founder_id) "
        "ORDER BY a.id DESC", (email, email)).fetchall()
    out = []
    for a in rows:
        fb = _founder_brief(a["founder_id"])
        f = conn.execute("SELECT stage, sector FROM founder WHERE id=?", (a["founder_id"],)).fetchone()
        rv = conn.execute("SELECT verdict, status FROM review WHERE founder_id=? ORDER BY round DESC LIMIT 1",
                          (a["founder_id"],)).fetchone()
        review_status = "not_started"
        if rv:
            if rv["verdict"] == "awarded" and rv["status"] == "recommended":
                review_status = "recommended"
            elif rv["verdict"] == "awarded" and rv["status"] == "issued":
                review_status = "mark_issued"
            elif rv["verdict"] == "awarded" and rv["status"] == "overridden":
                review_status = "overridden"
            else:
                review_status = "not_yet"
        out.append({"founder_id": a["founder_id"], "name": fb.get("name", ""),
                    "company": fb.get("company", ""), "stage": (f["stage"] if f else "") or "",
                    "sector": (f["sector"] if f else "") or "",
                    "assigned_at": a["assigned_at"] or "", "sla_due": a["sla_due"] or "",
                    "review_status": review_status})
    conn.close()
    return out

def reviewer_submit(reviewer_email, founder_id, verdict, gaps, note):
    """A reviewer submits a verdict — allowed ONLY for a founder actually assigned
    to them. Award becomes a recommendation awaiting countersign."""
    reviewer_email = (reviewer_email or "").strip().lower()
    if verdict not in ("not_yet", "awarded"):
        return {"error": "verdict must be 'not_yet' or 'awarded'.", "status": 400}
    conn = _db()
    a = conn.execute("SELECT * FROM review_assignment WHERE founder_id=? AND reviewer_email=? "
                     "ORDER BY id DESC LIMIT 1", (founder_id, reviewer_email)).fetchone()
    rname = None
    if a:
        rname = a["reviewer_name"]
    else:
        r = conn.execute("SELECT name FROM reviewer WHERE email=?", (reviewer_email,)).fetchone()
        rname = r["name"] if r else None
    if not a:
        conn.close()
        return {"error": "This founder isn't assigned to you. Ask TiE to assign it first.",
                "status": 403}
    conn.execute("UPDATE review_assignment SET status='reviewed' WHERE id=?", (a["id"],))
    conn.commit()
    conn.close()
    return review_add(founder_id, verdict, gaps, rname or reviewer_email, note)

def pending_verdicts(id_token):
    """Admin: awards a reviewer has recommended, awaiting countersign."""
    g = require_admin(id_token)
    if g.get("error"):
        return g
    conn = _db()
    rows = conn.execute(
        "SELECT id, founder_id, round, gaps_json, reviewer, note, created_at "
        "FROM review WHERE verdict='awarded' AND status='recommended' ORDER BY id DESC").fetchall()
    out = []
    for r in rows:
        fb = _founder_brief(r["founder_id"])
        out.append({"review_id": r["id"], "founder_id": r["founder_id"], "name": fb.get("name", ""),
                    "company": fb.get("company", ""), "round": r["round"],
                    "reviewer": r["reviewer"] or "", "note": r["note"] or "",
                    "gaps": json.loads(r["gaps_json"] or "[]"), "created_at": r["created_at"] or ""})
    conn.close()
    return {"pending": out}

def review_countersign(id_token, review_id, decision, note):
    """Admin: issue the mark (final) or override (withhold). Only now is the
    founder told about an award."""
    g = require_admin(id_token)
    if g.get("error"):
        return g
    if decision not in ("issue", "override"):
        return {"error": "decision must be 'issue' or 'override'.", "status": 400}
    conn = _db()
    row = conn.execute("SELECT * FROM review WHERE id=? AND verdict='awarded' AND status='recommended'",
                       (review_id,)).fetchone()
    if not row:
        conn.close()
        return {"error": "No pending recommendation found for that review.", "status": 404}
    fid = row["founder_id"]
    new_status = "issued" if decision == "issue" else "overridden"
    conn.execute("UPDATE review SET status=?, countersigned_by=?, countersigned_at=datetime('now') WHERE id=?",
                 (new_status, g["email"], review_id))
    conn.commit()
    conn.close()
    fb = _founder_brief(fid)
    if decision == "issue":
        if fb["email"]:
            notify("mark_awarded", [fb["email"]],
                   "You have earned the VentureReady mark",
                   "Congratulations " + (fb["name"] or "") +
                   " — TiE has countersigned your review and awarded the VentureReady mark. "
                   "Screened investors can now discover your profile.", "high")
        notify("mark_awarded_admin", ["admins"],
               "VentureReady mark issued",
               g["email"] + " countersigned and issued the mark to " +
               (fb["name"] or ("founder #" + str(fid))) + ".", "normal")
    else:
        if fb["email"]:
            notify("mark_withheld", [fb["email"]],
                   "An update on your VentureReady review",
                   "After TiE's final review, the mark is being held for now" +
                   ((" — note: " + note) if note else "") +
                   ". Please see your feedback for what to strengthen before the next round.", "normal")
        notify("mark_withheld_admin", ["admins"],
               "Mark recommendation overridden",
               g["email"] + " overrode the recommended mark for " +
               (fb["name"] or ("founder #" + str(fid))) + ".", "normal")
    return {"ok": True, "pending": pending_verdicts(id_token).get("pending", [])}

# ---- Data room (diligence documents) ----
DATAROOM_DIR = os.path.join(HERE, "dataroom")

def dataroom_add(founder_id, item_key, filename, raw):
    """Save a diligence document and point the checklist item at it.
    One document per checklist item — re-uploading replaces the previous file's row."""
    if not os.path.isdir(DATAROOM_DIR):
        os.makedirs(DATAROOM_DIR)
    safe = "".join(c for c in (filename or "doc.pdf") if c.isalnum() or c in "._- ")
    stored = "f%s_%s_%d_%s" % (founder_id or 0, (item_key or "item")[:24], int(time.time()), safe)
    with open(os.path.join(DATAROOM_DIR, stored), "wb") as fh:
        fh.write(raw)
    conn = _db()
    conn.execute("DELETE FROM dataroom_doc WHERE founder_id=? AND item_key=?", (founder_id, item_key))
    cur = conn.execute("INSERT INTO dataroom_doc(founder_id,item_key,filename,stored_path,size,uploaded_at) "
                       "VALUES(?,?,?,?,?,datetime('now'))",
                       (founder_id, item_key, filename, stored, len(raw)))
    did = cur.lastrowid
    conn.commit()
    conn.close()
    return did

def dataroom_list(founder_id):
    """Everything this founder has supplied, plus who has opened what."""
    conn = _db()
    docs = [dict(r) for r in conn.execute(
        "SELECT id,item_key,filename,size,uploaded_at FROM dataroom_doc WHERE founder_id=? "
        "ORDER BY item_key", (founder_id,)).fetchall()]
    views = [dict(r) for r in conn.execute(
        "SELECT item_key,viewer,viewed_at FROM dataroom_view WHERE founder_id=? "
        "ORDER BY id DESC LIMIT 25", (founder_id,)).fetchall()]
    conn.close()
    return {"docs": docs, "views": views}

def dataroom_view_log(founder_id, doc_id, item_key, viewer):
    conn = _db()
    conn.execute("INSERT INTO dataroom_view(founder_id,doc_id,item_key,viewer,viewed_at) "
                 "VALUES(?,?,?,?,datetime('now'))", (founder_id, doc_id, item_key, viewer or "Unknown"))
    conn.commit()
    conn.close()
    return {"ok": True}

def diagnostic_find_reusable(founder_id, filename, size):
    """Return a prior AI read for the same founder + identical deck (same filename and
    byte-size) so we never pay the API twice for the same deck. None if nothing to reuse."""
    conn = _db()
    row = conn.execute(
        "SELECT dg.live, dg.summary, dg.findings_json "
        "FROM diagnostic dg JOIN deck k ON k.id = dg.deck_id "
        "WHERE dg.founder_id=? AND k.filename=? AND k.size=? "
        "ORDER BY dg.id DESC LIMIT 1", (founder_id, filename, size)).fetchone()
    conn.close()
    if not row:
        return None
    return {"live": bool(row["live"]), "summary": row["summary"],
            "findings": json.loads(row["findings_json"] or "[]"), "reused": True}

def diagnostic_save(founder_id, deck_id, result):
    conn = _db()
    conn.execute("INSERT INTO diagnostic(founder_id,deck_id,live,summary,findings_json,created_at) "
                 "VALUES(?,?,?,?,?,datetime('now'))",
                 (founder_id, deck_id, 1 if result.get("live") else 0,
                  result.get("summary", ""), json.dumps(result.get("findings", []))))
    conn.commit()
    conn.close()

# ---- load .env (private config: Zoho keys + optional Anthropic key) ----
ENV = {}
try:
    for line in open(os.path.join(HERE, ".env")):
        line = line.strip()
        if line and "=" in line and not line.startswith("#"):
            k, v = line.split("=", 1)
            ENV[k.strip()] = v.strip()
except FileNotFoundError:
    print("WARNING: .env not found — Zoho check will not work.")

ZOHO_CLIENT_ID = ENV.get("ZOHO_CLIENT_ID", "")
ZOHO_CLIENT_SECRET = ENV.get("ZOHO_CLIENT_SECRET", "")
ANTHROPIC_API_KEY = ENV.get("ANTHROPIC_API_KEY", "")
# Google sign-in ("Continue with Google"). This is the PUBLIC OAuth Client ID —
# safe to expose to the browser (that's how Google's button works). There is NO
# client secret here: the browser gets a signed ID token and the server only
# verifies it, so no secret is needed. Blank = feature off (button falls back to demo).
GOOGLE_CLIENT_ID = ENV.get("GOOGLE_CLIENT_ID", "")

# Where founders go to actually BUY a TiE Bangalore membership. Membership money
# is not taken in this app — Zoho owns the membership record, so we hand off to
# it and then re-check the result. Set TIE_MEMBERSHIP_URL in .env; if it's blank
# the front-end hides the join buttons rather than showing a broken link.
TIE_MEMBERSHIP_URL = ENV.get("TIE_MEMBERSHIP_URL", "").strip()

# ---- Admin / Team-Portal access allow-list ----
# Only these Google-verified emails may enter the admin portal. Access is an
# explicit allow-list (NOT the whole @tiebangalore.org domain).
#
# SUPER-ADMINS are the trust anchor: they can add/remove admins and cannot be
# removed in-app. Set them (comma-separated) via ADMIN_SUPER_EMAILS in .env;
# default = the two TiE Bangalore owners.
# ADMIN_EMAILS from .env is only a FIRST-RUN SEED. The live allow-list lives in
# the admin_allowlist table (see db_init) so a super-admin can add/remove
# admins from inside the portal without editing files or restarting the server.
_super_raw = ENV.get("ADMIN_SUPER_EMAILS") or ENV.get("ADMIN_SUPER_EMAIL") or \
    "admin.blr@tiebangalore.org,chinmay@tiebangalore.org"
ADMIN_SUPER_EMAILS = set()
for _e in _super_raw.split(","):
    _e = _e.strip().lower()
    if _e:
        ADMIN_SUPER_EMAILS.add(_e)

def is_super_admin(email):
    return (email or "").strip().lower() in ADMIN_SUPER_EMAILS

def admin_role_for(email):
    email = (email or "").strip().lower()
    if not email:
        return None
    if is_super_admin(email):
        return "super-admin"
    conn = _db()
    row = conn.execute("SELECT email FROM admin_allowlist WHERE email=?", (email,)).fetchone()
    conn.close()
    return "admin" if row else None

def admin_list():
    # Super-admins sort to the top; everyone else alphabetically.
    conn = _db()
    rows = conn.execute(
        "SELECT email, added_by, added_at FROM admin_allowlist ORDER BY email ASC").fetchall()
    conn.close()
    mapped = [{"email": r["email"],
             "role": "super-admin" if is_super_admin(r["email"]) else "admin",
             "added_by": r["added_by"] or "",
             "added_at": r["added_at"] or ""} for r in rows]
    mapped.sort(key=lambda a: (0 if a["role"] == "super-admin" else 1, a["email"]))
    return mapped

def require_super_admin(id_token):
    # Verify a Google credential and require it to be a SUPER-admin. Guards the
    # admin-management endpoints: there are no sessions yet, so the caller proves
    # who they are by sending their Google ID token with each request.
    v = verify_google_token(id_token)
    if v.get("error"):
        return {"error": v["error"], "status": 401}
    if not is_super_admin(v["email"]):
        return {"error": "Only a super-admin can manage admin access.", "status": 403}
    return {"email": v["email"]}

def admin_add(id_token, email):
    g = require_super_admin(id_token)
    if g.get("error"):
        return g
    email = (email or "").strip().lower()
    if not re.match(r"^[^@\s]+@[^@\s]+\.[^@\s]+$", email):
        return {"error": "Please enter a valid email address.", "status": 400}
    conn = _db()
    conn.execute("INSERT OR IGNORE INTO admin_allowlist(email, added_by, added_at) "
                 "VALUES(?,?,datetime('now'))", (email, g["email"]))
    conn.commit()
    conn.close()
    notify("admin_added", [email] + sorted(ADMIN_SUPER_EMAILS),
           "You have been added as a VentureReady admin",
           g["email"] + " added " + email +
           " to the VentureReady Team Portal admin allow-list.", "normal")
    return {"admins": admin_list()}

def admin_remove(id_token, email):
    g = require_super_admin(id_token)
    if g.get("error"):
        return g
    email = (email or "").strip().lower()
    if is_super_admin(email):
        return {"error": "A super-admin can’t be removed.", "status": 400}
    conn = _db()
    conn.execute("DELETE FROM admin_allowlist WHERE email=?", (email,))
    conn.commit()
    conn.close()
    notify("admin_removed", sorted(ADMIN_SUPER_EMAILS),
           "An admin was removed from VentureReady",
           g["email"] + " removed " + email +
           " from the VentureReady Team Portal admin allow-list.", "normal")
    return {"admins": admin_list()}

# ---- Notifications & email outbox -----------------------------------------
# Every important event is RECORDED here (a notification plus one email_outbox
# row per recipient). Whether a recorded message is actually EMAILED depends on
# MAIL_ENABLED, which is true only when a full SMTP configuration is present in
# .env. Until then the platform runs in "record-only" mode: messages are
# captured and shown on the admin Notifications screen, but nothing is sent.
# This keeps the app safe to run before a mail provider is wired up.
SMTP_HOST = ENV.get("SMTP_HOST", "").strip()
SMTP_PORT = int((ENV.get("SMTP_PORT", "") or "587").strip() or "587")
SMTP_USER = ENV.get("SMTP_USER", "").strip()
SMTP_PASS = ENV.get("SMTP_PASS", "").strip()
MAIL_FROM = ENV.get("MAIL_FROM", "").strip() or SMTP_USER
MAIL_FROM_NAME = (ENV.get("MAIL_FROM_NAME", "") or "VentureReady").strip()
MAIL_ENABLED = bool(SMTP_HOST and SMTP_USER and SMTP_PASS and MAIL_FROM)

def _admin_emails():
    conn = _db()
    rows = conn.execute("SELECT email FROM admin_allowlist ORDER BY email ASC").fetchall()
    conn.close()
    return [r["email"] for r in rows]

def _founder_brief(founder_id):
    # Name/email/company for composing a notification about a founder.
    conn = _db()
    row = conn.execute("SELECT name, email, company FROM founder WHERE id=?",
                       (founder_id,)).fetchone()
    conn.close()
    if not row:
        return {"name": "", "email": "", "company": ""}
    return {"name": row["name"] or "", "email": row["email"] or "",
            "company": row["company"] or ""}

def _expand_recipients(recipients):
    # Accept a list of email addresses and/or the token "admins" (= everyone on
    # the allow-list). Returns a de-duplicated list with blanks removed.
    out = []
    for r in (recipients or []):
        r = (r or "").strip()
        if not r:
            continue
        if r.lower() == "admins":
            out.extend(_admin_emails())
        else:
            out.append(r)
    seen, uniq = set(), []
    for e in out:
        k = e.strip().lower()
        if k and k not in seen:
            seen.add(k)
            uniq.append(e)
    return uniq

def notify(event, recipients, subject, body, urgency="normal"):
    # Record one notification + one pending email per recipient. Actual sending
    # happens later (Stage 3) and only when MAIL_ENABLED; for now the rows simply
    # wait in the outbox. Returns how many recipients were recorded.
    conn = _db()
    n = 0
    for to in _expand_recipients(recipients):
        cur = conn.execute(
            "INSERT INTO notification(event, recipient, subject, body, urgency, created_at) "
            "VALUES(?,?,?,?,?,datetime('now'))", (event, to, subject, body, urgency))
        conn.execute(
            "INSERT INTO email_outbox(notification_id, to_email, subject, body, status, created_at) "
            "VALUES(?,?,?,?,'pending',datetime('now'))", (cur.lastrowid, to, subject, body))
        n += 1
    conn.commit()
    conn.close()
    return n

def require_admin(id_token):
    # Any admin (super-admin OR regular) may view the notifications feed.
    v = verify_google_token(id_token)
    if v.get("error"):
        return {"error": v["error"], "status": 401}
    role = admin_role_for(v["email"])
    if not role:
        return {"error": "Admin access required.", "status": 403}
    return {"email": v["email"], "role": role}

def notifications_recent(limit=100):
    conn = _db()
    rows = conn.execute(
        "SELECT n.id, n.event, n.recipient, n.subject, n.body, n.urgency, n.created_at, "
        "o.status FROM notification n LEFT JOIN email_outbox o ON o.notification_id = n.id "
        "ORDER BY n.id DESC LIMIT ?", (limit,)).fetchall()
    conn.close()
    return [{"id": r["id"], "event": r["event"], "recipient": r["recipient"],
             "subject": r["subject"], "body": r["body"],
             "urgency": r["urgency"] or "normal", "created_at": r["created_at"] or "",
             "status": r["status"] or "pending"} for r in rows]

# ---- Payments (real records + receipt notification) -----------------------
def payment_record(d):
    # Record a payment. No card details are ever stored — only what was bought.
    # Fires a receipt to the payer and a heads-up to admins.
    payer_email = (d.get("payer_email") or "").strip().lower()
    payer_name = (d.get("payer_name") or "").strip()
    founder_id = d.get("founder_id")
    item = (d.get("item") or "").strip()
    amount = (d.get("amount") or "").strip()
    period = (d.get("period") or "").strip()
    # Fall back to the logged-in founder's email if the caller didn't send one.
    if not payer_email and founder_id:
        fb = _founder_brief(founder_id)
        payer_email = (fb.get("email") or "").strip().lower()
        payer_name = payer_name or fb.get("name") or ""
    # Don't record an empty shell of a payment — we need at least what was
    # bought, or someone to attribute it to.
    if not item and not payer_email:
        return {"error": "A payment needs at least an item or a payer.", "status": 400}
    conn = _db()
    cur = conn.execute(
        "INSERT INTO payment(payer_email, payer_name, founder_id, item, amount, period, "
        "status, created_at) VALUES(?,?,?,?,?,?, 'paid', datetime('now'))",
        (payer_email, payer_name, founder_id, item, amount, period))
    pid = cur.lastrowid
    conn.commit()
    row = dict(conn.execute("SELECT * FROM payment WHERE id=?", (pid,)).fetchone())
    conn.close()
    amt_str = (amount + (" " + period if period else "")).strip()
    if payer_email:
        notify("payment_receipt", [payer_email],
               "Your VentureReady payment receipt",
               "We've received your payment of " + (amt_str or "your fee") + " for " +
               (item or "VentureReady access") + ". Thank you — your access is now active.",
               "normal")
    notify("payment_received", ["admins"],
           "Payment received: " + (item or "VentureReady access"),
           (payer_name or payer_email or "A user") + " paid " + (amt_str or "a fee") +
           " for " + (item or "VentureReady access") + ".", "normal")
    return {"payment": row}

# ---- Reviewer assignment (real records + reviewer notification) -----------
def review_queue():
    # Founders who have submitted at least one deck are candidates for review.
    # For each, show their most recent reviewer assignment (if any).
    conn = _db()
    rows = conn.execute(
        "SELECT f.id, f.name, f.company, f.stage, f.sector, "
        "  (SELECT MAX(uploaded_at) FROM deck d WHERE d.founder_id=f.id) AS last_deck, "
        "  (SELECT COUNT(*) FROM deck d WHERE d.founder_id=f.id) AS deck_count "
        "FROM founder f "
        "WHERE (SELECT COUNT(*) FROM deck d WHERE d.founder_id=f.id) > 0 "
        "ORDER BY last_deck DESC").fetchall()
    out = []
    for r in rows:
        a = conn.execute(
            "SELECT reviewer_email, reviewer_name, status, created_at FROM review_assignment "
            "WHERE founder_id=? ORDER BY id DESC LIMIT 1", (r["id"],)).fetchone()
        out.append({"founder_id": r["id"], "name": r["name"] or "", "company": r["company"] or "",
                    "stage": r["stage"] or "", "sector": r["sector"] or "",
                    "last_deck": r["last_deck"] or "", "deck_count": r["deck_count"] or 0,
                    "reviewer_email": (a["reviewer_email"] if a else ""),
                    "reviewer_name": (a["reviewer_name"] if a else ""),
                    "assignment_status": (a["status"] if a else "")})
    conn.close()
    return out

def review_assign(id_token, founder_id, reviewer_email, reviewer_name):
    g = require_admin(id_token)
    if g.get("error"):
        return g
    reviewer_email = (reviewer_email or "").strip().lower()
    if not re.match(r"^[^@\s]+@[^@\s]+\.[^@\s]+$", reviewer_email):
        return {"error": "Please enter a valid reviewer email address.", "status": 400}
    fb = _founder_brief(founder_id)
    if not fb.get("email") and not fb.get("name"):
        return {"error": "Unknown founder.", "status": 404}
    reviewer_name = (reviewer_name or "").strip() or reviewer_email.split("@")[0]
    prior = None
    conn = _db()
    prior = conn.execute("SELECT COUNT(*) c FROM review_assignment WHERE founder_id=?",
                         (founder_id,)).fetchone()["c"]
    status = "reassigned" if prior else "assigned"
    conn.execute(
        "INSERT INTO review_assignment(founder_id, founder_name, reviewer_email, reviewer_name, "
        "assigned_by, status, sla_due, created_at) VALUES(?,?,?,?,?,?, date('now','+5 day'), datetime('now'))",
        (founder_id, fb.get("name") or "", reviewer_email, reviewer_name, g["email"], status))
    conn.commit()
    conn.close()
    startup = (fb.get("company") or fb.get("name") or "a founder")
    notify("review_assigned", [reviewer_email],
           "You've been assigned a VentureReady review",
           "TiE Bangalore has assigned you to review " + startup + "'s submission. "
           "Target turnaround is 3–5 business days. Please open the Team Portal to begin.", "normal")
    notify("review_assigned_admin", ["admins"],
           "Review assigned: " + startup,
           g["email"] + " assigned " + startup + "'s review to " + reviewer_name +
           " (" + reviewer_email + ").", "normal")
    return {"ok": True, "queue": review_queue()}

# ---- Investors (real screened-access records) -----------------------------
def _investor_row(r):
    if r is None:
        return None
    return {"id": r["id"], "name": r["name"] or "", "email": r["email"] or "",
            "firm": r["firm"] or "", "tie_status": r["tie_status"] or "",
            "cheque_size": r["cheque_size"] or "", "focus": r["focus"] or "",
            "track_record": r["track_record"] or "",
            "recent_investments": r["recent_investments"] or "",
            "status": r["status"] or "pending",
            "parent_investor_id": r["parent_investor_id"],
            "decided_by": r["decided_by"] or "", "decided_at": r["decided_at"] or "",
            "created_at": r["created_at"] or ""}

def investor_apply(d):
    # A screened-access application from the public investor form. Creates (or
    # refreshes) a 'pending' investor record and tells the admins to review it.
    email = (d.get("email") or "").strip().lower()
    name = (d.get("name") or "").strip()
    if not re.match(r"^[^@\s]+@[^@\s]+\.[^@\s]+$", email):
        return {"error": "Please enter a valid email address.", "status": 400}
    if not name:
        return {"error": "Please enter your full name.", "status": 400}
    fields = {"firm": (d.get("firm") or "").strip(),
              "tie_status": (d.get("tie_status") or "").strip(),
              "cheque_size": (d.get("cheque_size") or "").strip(),
              "focus": (d.get("focus") or "").strip(),
              "track_record": (d.get("track_record") or "").strip(),
              "recent_investments": (d.get("recent_investments") or "").strip()}
    conn = _db()
    existing = conn.execute("SELECT * FROM investor WHERE email=?", (email,)).fetchone()
    if existing and (existing["status"] or "") == "approved":
        conn.close()
        return {"error": "That email already has approved investor access. Please sign in instead.",
                "status": 400}
    if existing:
        conn.execute("UPDATE investor SET name=?, firm=?, tie_status=?, cheque_size=?, focus=?, "
                     "track_record=?, recent_investments=?, status='pending' WHERE id=?",
                     (name, fields["firm"], fields["tie_status"], fields["cheque_size"],
                      fields["focus"], fields["track_record"], fields["recent_investments"],
                      existing["id"]))
        iid = existing["id"]
    else:
        cur = conn.execute(
            "INSERT INTO investor(name, email, firm, tie_status, cheque_size, focus, track_record, "
            "recent_investments, status, created_at) VALUES(?,?,?,?,?,?,?,?, 'pending', datetime('now'))",
            (name, email, fields["firm"], fields["tie_status"], fields["cheque_size"],
             fields["focus"], fields["track_record"], fields["recent_investments"]))
        iid = cur.lastrowid
    conn.commit()
    row = conn.execute("SELECT * FROM investor WHERE id=?", (iid,)).fetchone()
    conn.close()
    notify("investor_application", ["admins"],
           "New investor application: " + (fields["firm"] or name),
           name + " (" + email + ")" + (" from " + fields["firm"] if fields["firm"] else "") +
           " applied for screened deal-flow access. TiE status: " +
           (fields["tie_status"] or "not stated") + ". Review them in the Team Portal.", "normal")
    return {"investor": _investor_row(row)}

def investor_list(id_token):
    g = require_admin(id_token)
    if g.get("error"):
        return g
    conn = _db()
    rows = conn.execute("SELECT * FROM investor ORDER BY id DESC").fetchall()
    conn.close()
    return {"investors": [_investor_row(r) for r in rows]}

def investor_decision(id_token, investor_id, decision, reason=""):
    # Admin approves or rejects a screened-access application; the investor is told.
    g = require_admin(id_token)
    if g.get("error"):
        return g
    if decision not in ("approved", "rejected"):
        return {"error": "decision must be 'approved' or 'rejected'.", "status": 400}
    conn = _db()
    row = conn.execute("SELECT * FROM investor WHERE id=?", (investor_id,)).fetchone()
    if not row:
        conn.close()
        return {"error": "Unknown investor.", "status": 404}
    conn.execute("UPDATE investor SET status=?, decided_by=?, decided_at=datetime('now') WHERE id=?",
                 (decision, g["email"], investor_id))
    conn.commit()
    conn.close()
    who = row["name"] or row["email"]
    if decision == "approved":
        notify("investor_approved", [row["email"]],
               "Your VentureReady investor access is approved",
               "Good news — TiE Bangalore has approved your application for screened deal-flow "
               "access. The next step is to sign your NDA, after which you can view "
               "VentureReady-marked founders.", "normal")
    else:
        notify("investor_rejected", [row["email"]],
               "About your VentureReady investor application",
               "Thank you for applying for screened deal-flow access. After review, TiE Bangalore "
               "is not able to approve your application at this time." +
               ((" Note: " + reason) if reason else ""), "normal")
    notify("investor_decision_admin", ["admins"],
           "Investor " + decision + ": " + who,
           g["email"] + " " + decision + " " + who + " (" + row["email"] + ").", "normal")
    return {"ok": True, "investors": investor_list(id_token).get("investors", [])}

INVESTOR_TEAM_SEATS = 3

def investor_invite(d):
    # An approved investor adds up to 3 colleagues from their firm. Each becomes
    # its own 'invited' investor row (each must sign their own NDA), and each is
    # notified individually — the inviter's signature does NOT cover them.
    inviter_email = (d.get("inviter_email") or "").strip().lower()
    users = d.get("users") or []
    conn = _db()
    inviter = conn.execute("SELECT * FROM investor WHERE email=?", (inviter_email,)).fetchone()
    if not inviter:
        conn.close()
        return {"error": "We couldn't find your investor record. Please apply first.", "status": 404}
    if (inviter["status"] or "") != "approved":
        conn.close()
        return {"error": "Your investor access isn't approved yet, so you can't add team members.",
                "status": 403}
    parent_id = inviter["parent_investor_id"] or inviter["id"]
    used = conn.execute("SELECT COUNT(*) c FROM investor WHERE parent_investor_id=?",
                        (parent_id,)).fetchone()["c"]
    added, skipped = [], []
    for u in users:
        email = ((u or {}).get("email") or "").strip().lower()
        name = ((u or {}).get("name") or "").strip()
        if not email:
            continue
        if not re.match(r"^[^@\s]+@[^@\s]+\.[^@\s]+$", email):
            skipped.append({"email": email, "why": "not a valid email address"})
            continue
        if used + len(added) >= INVESTOR_TEAM_SEATS:
            skipped.append({"email": email, "why": "no team seats left (limit %d)" % INVESTOR_TEAM_SEATS})
            continue
        if conn.execute("SELECT 1 FROM investor WHERE email=?", (email,)).fetchone():
            skipped.append({"email": email, "why": "already has a VentureReady investor record"})
            continue
        conn.execute(
            "INSERT INTO investor(name, email, firm, tie_status, status, parent_investor_id, created_at) "
            "VALUES(?,?,?,?, 'invited', ?, datetime('now'))",
            (name or email.split("@")[0], email, inviter["firm"] or "", inviter["tie_status"] or "",
             parent_id))
        added.append({"email": email, "name": name or email.split("@")[0]})
    conn.commit()
    conn.close()
    firm = inviter["firm"] or "your firm"
    for a in added:
        notify("investor_user_invited", [a["email"]],
               "You've been invited to VentureReady deal flow",
               (inviter["name"] or inviter_email) + " has added you to " + firm +
               "'s VentureReady deal-flow access. You must verify your email and sign your own "
               "NDA before you can view any founder material — their signature does not cover you.",
               "normal")
    if added:
        notify("investor_users_added", [inviter_email],
               "Your firm's team invitations were sent",
               "You added " + str(len(added)) + " colleague(s) to " + firm +
               "'s VentureReady access: " + ", ".join(a["email"] for a in added) +
               ". Each must sign their own NDA.", "normal")
        notify("investor_users_added_admin", ["admins"],
               "Investor team seats used: " + firm,
               (inviter["name"] or inviter_email) + " added " + str(len(added)) +
               " team member(s) to " + firm + ": " + ", ".join(a["email"] for a in added) + ".",
               "normal")
    return {"added": added, "skipped": skipped,
            "seats_left": max(0, INVESTOR_TEAM_SEATS - (used + len(added)))}

# ---- Membership hand-off (money is taken in Zoho, never in this app) -------
def membership_handoff(d):
    # Records that a founder was sent to Zoho to join, so admins can follow up
    # on people who start but never finish. This is NOT a payment record — the
    # app never sees the money and never claims they paid.
    founder_id = d.get("founder_id")
    fb = _founder_brief(founder_id) if founder_id else {}
    email = (d.get("email") or fb.get("email") or "").strip().lower()
    name = (d.get("name") or fb.get("name") or "").strip()
    tier = (d.get("tier") or "Associate Member").strip()
    # Only tell admins when we know WHO set out to join — an alert naming nobody
    # can't be followed up on, so it would just be noise in the feed.
    if not (email or name):
        return {"ok": True, "recorded": False}
    notify("membership_handoff", ["admins"],
           "Founder started joining TiE: " + (name or email),
           (name or email) + " opened the TiE Bangalore membership form to join as " + tier +
           ". They're recognised as a member as soon as they finish; the app re-checks when they "
           "come back, so follow up if it never appears.", "normal")
    return {"ok": True, "recorded": True}

def membership_recheck(d):
    # After the founder says they've joined, ask Zoho again and cache the answer.
    # Zoho stays the source of truth; the app never marks anyone a member itself.
    founder_id = d.get("founder_id")
    if not founder_id:
        return {"error": "Please sign in first so we know whose membership to check.",
                "status": 400}
    row = refresh_membership(founder_id)
    if not row:
        return {"error": "We couldn't find your account.", "status": 404}
    pub = _public_founder(row) or {}
    if pub.get("is_member"):
        notify("membership_confirmed", [pub.get("email") or ""],
               "Your TiE Bangalore membership is confirmed",
               "Thanks " + (pub.get("name") or "") + " — we've confirmed your TiE Bangalore "
               "membership with our records. Your VentureReady expert review is now included.",
               "normal")
        notify("membership_confirmed_admin", ["admins"],
               "New TiE member confirmed: " + (pub.get("name") or pub.get("email") or ""),
               (pub.get("name") or pub.get("email") or "A founder") +
               " completed TiE Bangalore membership and it is now confirmed in Zoho.", "normal")
    return {"founder": pub}

# ---- Meeting requests (TiE-facilitated introductions) ----------------------
def meeting_request(d):
    # An investor asks TiE to introduce them to a founder. TiE facilitates every
    # intro, so this notifies the admins — never the founder directly.
    founder_id = d.get("founder_id")
    investor_email = (d.get("investor_email") or "").strip().lower()
    investor_name = (d.get("investor_name") or "").strip()
    fb = _founder_brief(founder_id) if founder_id else {}
    founder_name = (d.get("founder_name") or fb.get("name") or "").strip()
    startup = (d.get("startup") or fb.get("company") or "").strip()
    if not founder_name and not startup:
        return {"error": "We need to know which founder you'd like to meet.", "status": 400}
    investor_id = None
    conn = _db()
    if investor_email:
        inv = conn.execute("SELECT * FROM investor WHERE email=?", (investor_email,)).fetchone()
        if inv:
            investor_id = inv["id"]
            investor_name = investor_name or (inv["name"] or "")
    cur = conn.execute(
        "INSERT INTO meeting(founder_id, founder_name, startup, investor_id, investor_name, "
        "investor_email, status, requested_at) VALUES(?,?,?,?,?,?, 'requested', datetime('now'))",
        (founder_id, founder_name, startup, investor_id, investor_name, investor_email))
    mid = cur.lastrowid
    conn.commit()
    row = dict(conn.execute("SELECT * FROM meeting WHERE id=?", (mid,)).fetchone())
    conn.close()
    notify("meeting_requested", ["admins"],
           "Meeting request: " + (investor_name or investor_email or "An investor") + " → " +
           (startup or founder_name),
           (investor_name or investor_email or "An investor") + " has requested a TiE-facilitated "
           "introduction to " + (founder_name or "a founder") +
           (" (" + startup + ")" if startup else "") + ". Approve and send the intro from the "
           "Team Portal.", "normal")
    return {"meeting": row}

def meeting_list(id_token):
    g = require_admin(id_token)
    if g.get("error"):
        return g
    conn = _db()
    rows = [dict(r) for r in conn.execute("SELECT * FROM meeting ORDER BY id DESC").fetchall()]
    conn.close()
    return {"meetings": rows}

def meeting_intro_sent(id_token, meeting_id):
    # Admin confirms the introduction has been made: both sides are told.
    g = require_admin(id_token)
    if g.get("error"):
        return g
    conn = _db()
    row = conn.execute("SELECT * FROM meeting WHERE id=?", (meeting_id,)).fetchone()
    if not row:
        conn.close()
        return {"error": "Unknown meeting request.", "status": 404}
    conn.execute("UPDATE meeting SET status='intro_sent', intro_sent_at=datetime('now') WHERE id=?",
                 (meeting_id,))
    conn.commit()
    conn.close()
    fb = _founder_brief(row["founder_id"]) if row["founder_id"] else {}
    founder_email = fb.get("email") or ""
    startup = row["startup"] or fb.get("company") or ""
    inv_who = row["investor_name"] or row["investor_email"] or "an investor"
    if founder_email:
        notify("meeting_intro_sent_founder", [founder_email],
               "An investor introduction has been made",
               "TiE Bangalore has introduced you to " + inv_who +
               ", who asked to meet after seeing your VentureReady profile. Look out for the "
               "introduction email and reply directly to arrange a time.", "normal")
    if row["investor_email"]:
        notify("meeting_intro_sent_investor", [row["investor_email"]],
               "Your introduction has been made",
               "TiE Bangalore has introduced you to " + (row["founder_name"] or "the founder") +
               (" (" + startup + ")" if startup else "") +
               ". You can now correspond directly to arrange a meeting.", "normal")
    notify("meeting_intro_sent_admin", ["admins"],
           "Intro sent: " + inv_who + " → " + (startup or row["founder_name"] or "founder"),
           g["email"] + " marked the introduction between " + inv_who + " and " +
           (row["founder_name"] or "the founder") + " as sent.", "normal")
    return {"ok": True, "meetings": meeting_list(id_token).get("meetings", [])}

# ---- Deals & impact (attributable investment tracking) --------------------
# The funnel a founder<->investor relationship walks. 'passed' is a dead end
# (deal died); the rest are forward progress ending in 'closed' (an investment).
DEAL_STAGES = ["introduced", "met", "diligence", "term_sheet", "closed", "passed"]
# Progress order for "reached at least X" counting (excludes the dead-end).
DEAL_PROGRESS = ["introduced", "met", "diligence", "term_sheet", "closed"]

def _deal_confidence(founder_confirmed, investor_confirmed, tie_verified):
    # How much weight an outcome can carry. TiE's own check outranks self-report;
    # both sides agreeing outranks one; one outranks none.
    if tie_verified:
        return "tie_verified"
    if founder_confirmed and investor_confirmed:
        return "both_confirmed"
    if founder_confirmed or investor_confirmed:
        return "one_confirmed"
    return "unconfirmed"

def _deal_row(r):
    if r is None:
        return None
    fc, ic, tv = bool(r["founder_confirmed"]), bool(r["investor_confirmed"]), bool(r["tie_verified"])
    return {"id": r["id"], "founder_id": r["founder_id"], "founder_name": r["founder_name"] or "",
            "startup": r["startup"] or "", "investor_id": r["investor_id"],
            "investor_name": r["investor_name"] or "", "investor_email": r["investor_email"] or "",
            "firm": r["firm"] or "", "stage": r["stage"] or "introduced",
            "amount_inr": r["amount_inr"], "amount_disclosed": bool(r["amount_disclosed"]),
            "founder_consent_public": bool(r["founder_consent_public"]),
            "investor_consent_public": bool(r["investor_consent_public"]),
            "founder_confirmed": fc, "investor_confirmed": ic, "tie_verified": tv,
            "confidence": _deal_confidence(fc, ic, tv),
            "created_at": r["created_at"] or "", "updated_at": r["updated_at"] or "",
            "closed_at": r["closed_at"] or "", "note": r["note"] or ""}

def deal_create(d):
    # Open a deal between a founder and an investor. Reuses an existing OPEN deal
    # for the same pair rather than creating duplicates.
    founder_id = d.get("founder_id")
    fb = _founder_brief(founder_id) if founder_id else {}
    founder_name = (d.get("founder_name") or fb.get("name") or "").strip()
    startup = (d.get("startup") or fb.get("company") or "").strip()
    investor_email = (d.get("investor_email") or "").strip().lower()
    investor_name = (d.get("investor_name") or "").strip()
    investor_id, firm = d.get("investor_id"), (d.get("firm") or "").strip()
    conn = _db()
    if investor_email:
        inv = conn.execute("SELECT * FROM investor WHERE email=?", (investor_email,)).fetchone()
        if inv:
            investor_id = inv["id"]; investor_name = investor_name or (inv["name"] or "")
            firm = firm or (inv["firm"] or "")
    if not founder_name and not startup:
        conn.close()
        return {"error": "A deal needs a founder.", "status": 400}
    if not (investor_email or investor_name):
        conn.close()
        return {"error": "A deal needs an investor.", "status": 400}
    # Reuse an open deal for the same pair (don't double-count the same relationship).
    existing = conn.execute(
        "SELECT * FROM deal WHERE founder_id IS ? AND investor_email=? "
        "AND stage NOT IN ('closed','passed') ORDER BY id DESC LIMIT 1",
        (founder_id, investor_email)).fetchone()
    if existing:
        row = _deal_row(existing); conn.close()
        return {"deal": row, "reused": True}
    cur = conn.execute(
        "INSERT INTO deal(founder_id, founder_name, startup, investor_id, investor_name, "
        "investor_email, firm, stage, source, meeting_id, created_at, updated_at) "
        "VALUES(?,?,?,?,?,?,?, 'introduced', ?, ?, datetime('now'), datetime('now'))",
        (founder_id, founder_name, startup, investor_id, investor_name, investor_email, firm,
         d.get("source") or "manual", d.get("meeting_id")))
    did = cur.lastrowid
    conn.commit()
    row = _deal_row(conn.execute("SELECT * FROM deal WHERE id=?", (did,)).fetchone())
    conn.close()
    notify("deal_created", ["admins"],
           "New deal tracked: " + (investor_name or investor_email or "investor") + " → " +
           (startup or founder_name),
           "A deal between " + (investor_name or investor_email or "an investor") + " and " +
           (startup or founder_name) + " is now being tracked toward outcome.", "normal")
    return {"deal": row, "reused": False}

def deal_advance(d):
    # Move a deal along, and/or record an amount + public-consent, attributed to
    # whoever acted. When a party changes the stage, the OTHER party's prior
    # confirmation is cleared so they must re-attest the new stage (keeps the
    # two-sided signal honest); an admin stage change clears both.
    deal_id = d.get("deal_id")
    actor = (d.get("actor") or "").strip().lower()   # 'founder' | 'investor' | 'admin'
    conn = _db()
    row = conn.execute("SELECT * FROM deal WHERE id=?", (deal_id,)).fetchone()
    if not row:
        conn.close()
        return {"error": "Unknown deal.", "status": 404}
    cur = dict(row)
    new_stage = (d.get("stage") or cur["stage"] or "introduced").strip()
    if new_stage not in DEAL_STAGES:
        conn.close()
        return {"error": "Unknown stage.", "status": 400}
    stage_changed = new_stage != (cur["stage"] or "introduced")
    fc, ic, tv = cur["founder_confirmed"], cur["investor_confirmed"], cur["tie_verified"]
    fcp, icp = cur["founder_consent_public"], cur["investor_consent_public"]
    # Attribution + consent for the acting party.
    consent = d.get("consent_public")
    if actor == "founder":
        fc = 1
        if consent is not None:
            fcp = 1 if consent else 0
    elif actor == "investor":
        ic = 1
        if consent is not None:
            icp = 1 if consent else 0
    # If the stage moved, stale the confirmations that no longer apply.
    if stage_changed:
        if actor == "founder":
            ic = 0; tv = 0
        elif actor == "investor":
            fc = 0; tv = 0
        else:  # admin recorded the stage — both parties should re-attest it
            fc = 0; ic = 0
    amount_inr, amount_disclosed = cur["amount_inr"], cur["amount_disclosed"]
    if d.get("amount_inr") not in (None, ""):
        try:
            amount_inr = int(d.get("amount_inr")); amount_disclosed = 1
        except (ValueError, TypeError):
            conn.close()
            return {"error": "Amount must be a whole number of rupees.", "status": 400}
    closed_at = cur["closed_at"]
    if new_stage == "closed" and not closed_at:
        closed_at = None  # set via SQL below
    conn.execute(
        "UPDATE deal SET stage=?, amount_inr=?, amount_disclosed=?, founder_consent_public=?, "
        "investor_consent_public=?, founder_confirmed=?, investor_confirmed=?, tie_verified=?, "
        "note=COALESCE(?, note), updated_at=datetime('now'), "
        "closed_at=CASE WHEN ?='closed' AND closed_at IS NULL THEN datetime('now') ELSE closed_at END "
        "WHERE id=?",
        (new_stage, amount_inr, amount_disclosed, fcp, icp, fc, ic, tv,
         d.get("note"), new_stage, deal_id))
    conn.commit()
    fresh = _deal_row(conn.execute("SELECT * FROM deal WHERE id=?", (deal_id,)).fetchone())
    conn.close()
    # Nudge the OTHER party to confirm, and flag closures to admins (that's impact).
    label = (fresh["startup"] or fresh["founder_name"]) + " ↔ " + (fresh["investor_name"] or fresh["investor_email"] or "investor")
    if new_stage == "closed":
        notify("deal_closed", ["admins"], "Deal closed: " + label,
               "A deal has been marked closed: " + label +
               (". Amount: ₹%s" % fresh["amount_inr"] if fresh["amount_disclosed"] and fresh["amount_inr"] else
                ". Amount undisclosed.") + " Confidence: " + fresh["confidence"] + ".", "high")
    elif stage_changed:
        # ask the party who did NOT just act to confirm
        if actor == "founder" and fresh["investor_email"]:
            notify("deal_update_confirm", [fresh["investor_email"]],
                   "Please confirm a deal update", "The status of your deal (" + label +
                   ") was updated to '" + new_stage.replace("_", " ") +
                   "'. Please confirm it so TiE's records stay accurate.", "normal")
        elif actor == "investor" and fresh["founder_id"]:
            fb2 = _founder_brief(fresh["founder_id"])
            if fb2.get("email"):
                notify("deal_update_confirm", [fb2["email"]],
                       "Please confirm a deal update", "The status of your deal (" + label +
                       ") was updated to '" + new_stage.replace("_", " ") +
                       "'. Please confirm it so TiE's records stay accurate.", "normal")
    return {"deal": fresh}

def deal_verify(id_token, deal_id, verified=True):
    g = require_admin(id_token)
    if g.get("error"):
        return g
    conn = _db()
    row = conn.execute("SELECT * FROM deal WHERE id=?", (deal_id,)).fetchone()
    if not row:
        conn.close()
        return {"error": "Unknown deal.", "status": 404}
    conn.execute("UPDATE deal SET tie_verified=?, updated_at=datetime('now') WHERE id=?",
                 (1 if verified else 0, deal_id))
    conn.commit()
    fresh = _deal_row(conn.execute("SELECT * FROM deal WHERE id=?", (deal_id,)).fetchone())
    conn.close()
    return {"deal": fresh}

def deals_for(founder_id=None, investor_email=None):
    # A single party's own deals (their private view).
    conn = _db()
    if founder_id:
        rows = conn.execute("SELECT * FROM deal WHERE founder_id=? ORDER BY id DESC", (founder_id,)).fetchall()
    elif investor_email:
        rows = conn.execute("SELECT * FROM deal WHERE investor_email=? ORDER BY id DESC",
                            ((investor_email or "").strip().lower(),)).fetchall()
    else:
        rows = []
    conn.close()
    return [_deal_row(r) for r in rows]

def deals_admin(id_token):
    g = require_admin(id_token)
    if g.get("error"):
        return g
    conn = _db()
    rows = [_deal_row(r) for r in conn.execute("SELECT * FROM deal ORDER BY id DESC").fetchall()]
    conn.close()
    # Funnel + capital, for TiE's own (full) view.
    by_stage = {s: 0 for s in DEAL_STAGES}
    conf = {"tie_verified": 0, "both_confirmed": 0, "one_confirmed": 0, "unconfirmed": 0}
    capital_disclosed = 0
    undisclosed_closed = 0
    closed = [r for r in rows if r["stage"] == "closed"]
    for r in rows:
        by_stage[r["stage"]] = by_stage.get(r["stage"], 0) + 1
    for r in closed:
        conf[r["confidence"]] = conf.get(r["confidence"], 0) + 1
        if r["amount_disclosed"] and r["amount_inr"]:
            capital_disclosed += r["amount_inr"]
        else:
            undisclosed_closed += 1
    stats = {"total": len(rows), "by_stage": by_stage, "closed": len(closed),
             "capital_disclosed_inr": capital_disclosed, "undisclosed_closed": undisclosed_closed,
             "confidence": conf,
             "founders": len({r["founder_id"] for r in rows if r["founder_id"]}),
             "investors": len({r["investor_email"] for r in rows if r["investor_email"]})}
    return {"deals": rows, "stats": stats}

def impact_public():
    # PUBLIC, showcase-safe aggregates ONLY. No names, no individual amounts.
    # The ₹ headline counts a closed deal ONLY when an amount was disclosed AND
    # both parties consented to public inclusion — everything else is a count.
    conn = _db()
    rows = [_deal_row(r) for r in conn.execute("SELECT * FROM deal").fetchall()]
    conn.close()
    def reached(stage_name):
        idx = DEAL_PROGRESS.index(stage_name)
        n = 0
        for r in rows:
            if r["stage"] in DEAL_PROGRESS and DEAL_PROGRESS.index(r["stage"]) >= idx:
                n += 1
        return n
    closed = [r for r in rows if r["stage"] == "closed"]
    capital_public = 0
    public_closings = 0
    for r in closed:
        if (r["amount_disclosed"] and r["amount_inr"]
                and r["founder_consent_public"] and r["investor_consent_public"]):
            capital_public += r["amount_inr"]
            public_closings += 1
    return {"introductions": len(rows), "meetings": reached("met"),
            "in_diligence_plus": reached("diligence"), "closed": len(closed),
            "startups_backed": len({r["founder_id"] for r in closed if r["founder_id"]}),
            "investors_active": len({r["investor_email"] for r in rows
                                     if r["investor_email"] and r["stage"] != "introduced"}),
            "capital_enabled_inr": capital_public,
            "capital_from_deals": public_closings,
            "closings_amount_undisclosed": len(closed) - public_closings}

CHAPTER = "TiE Bangalore"

# ---- Zoho token cache (token lives 1 hour) ----
_token = {"value": None, "expires": 0}

def zoho_token():
    if _token["value"] and time.time() < _token["expires"] - 60:
        return _token["value"]
    q = urllib.parse.urlencode({
        "grant_type": "client_credentials",
        "client_id": ZOHO_CLIENT_ID,
        "client_secret": ZOHO_CLIENT_SECRET,
        "scope": "ZohoCreator.report.READ",
    })
    req = urllib.request.Request("https://accounts.zoho.com/oauth/v2/token?" + q, method="POST")
    d = json.loads(urllib.request.urlopen(req, timeout=20).read())
    _token["value"] = d["access_token"]
    _token["expires"] = time.time() + int(d.get("expires_in", 3600))
    return _token["value"]

def verify_member(email):
    """Ask Zoho whether this email is a current TiE Bangalore member.

    Returns a dict. Three possible shapes:
      - Found & answered:  {member: bool, name, status, category}
      - Reached, no record: {member: False, reason: "not found", status: ""}
      - COULD NOT REACH Zoho: {member: False, unreachable: True, reason: ...}

    The `unreachable` flag matters: it lets callers tell "Zoho says they are not a
    member" apart from "we simply couldn't check right now", so we can keep the
    last CONFIRMED status instead of wrongly wiping someone's membership.
    """
    email = (email or "").strip()
    if not email:
        return {"member": False, "reason": "no email"}
    params = urllib.parse.urlencode({
        "field_config": "custom",
        "fields": "Email1,Name1,Status,Chapter_Name,Membership_Category1,Mobile",
        "Original_Email": email,
        "Chapter_Name": CHAPTER,
    })
    url = "https://creator.zoho.com/api/v2.1/tie_dev/chapters/report/Member_Details_Admin_View?" + params
    try:
        # zoho_token() can itself fail if Zoho's auth server is down — keep it
        # inside the try so an auth outage counts as "unreachable", not "not found".
        req = urllib.request.Request(url, headers={"Authorization": "Zoho-oauthtoken " + zoho_token()})
        d = json.loads(urllib.request.urlopen(req, timeout=20).read())
        rec = (d.get("data") or [{}])[0]
        status = rec.get("Status", "")
        return {
            "member": status == "Active",
            "name": (rec.get("Name1") or {}).get("zc_display_value", ""),
            "status": status,
            "category": (rec.get("Membership_Category1") or {}).get("zc_display_value", ""),
        }
    except urllib.error.HTTPError as e:
        # 5xx = Zoho itself is having trouble → treat as unreachable (keep last known).
        # 4xx (incl. code 9220 = no matching record) = a real answer: not a member.
        if e.code >= 500:
            return {"member": False, "unreachable": True, "reason": "zoho error %s" % e.code}
        return {"member": False, "reason": "not found", "status": ""}
    except Exception as e:
        # Network timeout, DNS failure, auth failure, bad JSON — we couldn't get a
        # trustworthy answer, so don't overwrite what we already know.
        return {"member": False, "unreachable": True, "reason": str(e)}

def refresh_membership(founder_id):
    """Re-check a founder's TiE membership with Zoho and cache the answer.

    Called at every sign-in (login and signup). Zoho stays the source of truth;
    the founder row just caches the latest answer so the profile shows it.

    If Zoho can't be reached we KEEP the last confirmed values and only flip
    membership_stale=1, so the UI can say "showing last confirmed status". A
    successful check clears the stale flag. Returns the fresh founder row.
    """
    conn = _db()
    row = conn.execute("SELECT * FROM founder WHERE id=?", (founder_id,)).fetchone()
    if not row:
        conn.close()
        return None
    res = verify_member(row["email"])
    if res.get("unreachable"):
        # Couldn't reach Zoho — leave the cached snapshot untouched, just flag it stale.
        conn.execute("UPDATE founder SET membership_stale=1 WHERE id=?", (founder_id,))
    else:
        conn.execute(
            "UPDATE founder SET is_member=?, membership_status=?, membership_category=?, "
            "membership_name=?, membership_checked_at=datetime('now'), membership_stale=0 "
            "WHERE id=?",
            (1 if res.get("member") else 0, res.get("status", ""),
             res.get("category", ""), res.get("name", ""), founder_id))
    conn.commit()
    out = conn.execute("SELECT * FROM founder WHERE id=?", (founder_id,)).fetchone()
    conn.close()
    return out

# ---- "Continue with Google" sign-in ----
# Google proves IDENTITY (who the person is). It does NOT prove TiE membership —
# that's still Zoho's job, run right after via refresh_membership(). The two gates
# stay separate: identity first, membership second.
def verify_google_token(id_token):
    """Verify the signed ID token the browser got from Google.

    We ask Google's own tokeninfo endpoint to validate it (so we're not
    hand-rolling JWT crypto), then still check audience/issuer/verified-email
    ourselves — never trust a token that wasn't minted for THIS app.

    NOTE FOR PRODUCTION: for higher volume, verify the token locally against
    Google's public keys with a vetted library instead of the tokeninfo endpoint.
    """
    if not GOOGLE_CLIENT_ID:
        return {"error": "Google sign-in is not configured on this server."}
    if not id_token:
        return {"error": "No Google credential was provided."}
    try:
        url = "https://oauth2.googleapis.com/tokeninfo?id_token=" + urllib.parse.quote(id_token)
        d = json.loads(urllib.request.urlopen(url, timeout=20).read())
    except urllib.error.HTTPError:
        return {"error": "Google could not verify that sign-in. Please try again."}
    except Exception:
        return {"error": "Couldn't reach Google to verify the sign-in. Please try again."}
    # aud must be OUR client ID — stops a token issued for another app being replayed here.
    if d.get("aud") != GOOGLE_CLIENT_ID:
        return {"error": "That sign-in was not issued for this app."}
    # iss must be Google.
    if d.get("iss") not in ("accounts.google.com", "https://accounts.google.com"):
        return {"error": "That sign-in did not come from Google."}
    # Token must not be expired.
    try:
        if d.get("exp") and time.time() > int(d["exp"]):
            return {"error": "That Google sign-in has expired. Please try again."}
    except (TypeError, ValueError):
        pass
    # Only accept a verified email — we key TiE membership off the email address.
    if str(d.get("email_verified")).lower() != "true" or not d.get("email"):
        return {"error": "Your Google email isn't verified, so we can't use it to sign in."}
    return {"sub": d.get("sub"), "email": (d.get("email") or "").strip().lower(), "name": d.get("name", "")}

def founder_google_upsert(sub, email, name):
    """Find the founder by Google ID first, then by email, else create a new record.
    Google accounts have no password — that's fine; password columns stay null."""
    conn = _db()
    email = (email or "").strip().lower()
    is_new = False
    row = conn.execute("SELECT * FROM founder WHERE google_sub=?", (sub,)).fetchone()
    if not row and email:
        row = conn.execute("SELECT * FROM founder WHERE email=?", (email,)).fetchone()
    if row:
        # Attach the Google ID (and fill a blank name) without disturbing anything else.
        conn.execute("UPDATE founder SET google_sub=?, name=COALESCE(NULLIF(name,''), ?) WHERE id=?",
                     (sub, name or "", row["id"]))
        fid = row["id"]
    else:
        cur = conn.execute(
            "INSERT INTO founder(name,email,google_sub,created_at) VALUES(?,?,?,datetime('now'))",
            (name or "New Founder", email, sub))
        fid = cur.lastrowid
        is_new = True
    conn.commit()
    conn.close()
    if is_new and email:
        notify("founder_welcome", [email],
               "Welcome to VentureReady",
               "Thanks for joining VentureReady, " + (name or "founder") +
               ". You can now run the free positioning diagnostic and submit your deck for expert review.",
               "normal")
    return fid

def founder_google_login(id_token):
    """Full flow: verify identity with Google → find/create the founder → re-check
    TiE membership with Zoho → return the public profile."""
    v = verify_google_token(id_token)
    if v.get("error"):
        return {"error": v["error"]}
    fid = founder_google_upsert(v["sub"], v["email"], v["name"])
    refreshed = refresh_membership(fid)
    if refreshed is None:
        conn = _db()
        refreshed = conn.execute("SELECT * FROM founder WHERE id=?", (fid,)).fetchone()
        conn.close()
    # Decide the role from the admin allow-list. Everyone else is a "founder".
    role = admin_role_for(v["email"]) or "founder"
    return {"founder": _public_founder(refreshed), "role": role}

# ---- Deck text extraction (reads the founder's actual uploaded deck) ----
# How many slide pictures we'll send the AI, and how big each may be.
# Vision costs more than text, so this path is a FALLBACK for image-only decks
# and is capped to keep a runaway deck from burning the budget.
MAX_DECK_IMAGES = 20
MAX_IMAGE_PX = 1400

def _shrink_to_jpeg(data, max_px=MAX_IMAGE_PX):
    """Normalise any embedded picture to a modest JPEG the API will accept."""
    from PIL import Image
    im = Image.open(io.BytesIO(data))
    if im.mode not in ("RGB", "L"):
        im = im.convert("RGB")
    w, h = im.size
    if max(w, h) > max_px:
        scale = max_px / float(max(w, h))
        im = im.resize((max(1, int(w * scale)), max(1, int(h * scale))), Image.LANCZOS)
    buf = io.BytesIO()
    im.save(buf, format="JPEG", quality=80, optimize=True)
    return buf.getvalue()

def extract_deck_images(filename, raw, max_images=MAX_DECK_IMAGES):
    """For image-only decks: pull the picture of each slide so the AI can LOOK at it.
    Takes the largest image per page/slide (that's the slide render), shrinks it,
    and returns base64 JPEGs ready for the API. Empty list if nothing usable."""
    name = (filename or "").lower()
    pics = []
    try:
        if name.endswith(".pdf"):
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(raw))
            for page in reader.pages[:max_images]:
                try:
                    imgs = list(page.images)
                except Exception:
                    imgs = []
                if not imgs:
                    continue
                biggest = max(imgs, key=lambda i: len(i.data))   # the slide render
                pics.append(biggest.data)
        elif name.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(raw))
            for slide in prs.slides:
                if len(pics) >= max_images:
                    break
                blobs = []
                for shape in slide.shapes:
                    try:
                        if getattr(shape, "image", None) is not None:
                            blobs.append(shape.image.blob)
                    except Exception:
                        pass
                if blobs:
                    pics.append(max(blobs, key=len))
    except Exception as e:
        print("deck image extraction failed:", e)
        return []

    out = []
    for data in pics:
        try:
            out.append({"media_type": "image/jpeg",
                        "b64": base64.b64encode(_shrink_to_jpeg(data)).decode()})
        except Exception as e:
            print("image convert skipped:", e)
    return out

def extract_deck_text(filename, raw):
    """Pull plain text out of an uploaded PDF or PPTX so the AI reads the real deck."""
    name = (filename or "").lower()
    try:
        if name.endswith(".pdf"):
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(raw))
            return "\n".join((page.extract_text() or "") for page in reader.pages).strip()
        if name.endswith(".pptx") or name.endswith(".ppt"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(raw))
            out = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        out.append(shape.text_frame.text)
            return "\n".join(out).strip()
    except Exception as e:
        print("deck extraction failed:", e)
        return ""
    # Unknown type: try decoding as plain text.
    try:
        return raw.decode("utf-8", "ignore").strip()
    except Exception:
        return ""


# ---- AI diagnostic ----
ANTI_PATTERNS = (
    "TAM inflation, vague ICP with no firmographic qualifier, no buyer trigger event, missing 'why now', "
    "feature-parity-as-moat, vanity metrics, wishful GTM, solution looking for a problem, internal "
    "contradictions, projection overconfidence, regulatory blindspots, and US-playbook assumptions in the Indian market"
)

# ---- Model auto-discovery ----
# Nobody has to know or maintain the AI model name. The app asks Anthropic which
# models the account can use and picks the strongest one, preferring Opus, then
# Sonnet. A manual override (ANTHROPIC_MODEL in .env) wins if ever needed.
_model_cache = {"value": None}

def get_model():
    if _model_cache["value"]:
        return _model_cache["value"]
    override = ENV.get("ANTHROPIC_MODEL", "").strip()
    if override:
        _model_cache["value"] = override
        return override
    fallback = "claude-sonnet-4-5"  # safe default if discovery ever fails
    try:
        req = urllib.request.Request(
            "https://api.anthropic.com/v1/models?limit=100", method="GET",
            headers={"x-api-key": ANTHROPIC_API_KEY, "anthropic-version": "2023-06-01"},
        )
        data = json.loads(urllib.request.urlopen(req, timeout=20).read()).get("data", [])
        ids = [m["id"] for m in data if "id" in m]
        # API returns newest first; keep that order and prefer Sonnet (best value for
        # this diagnostic), then fall back to Opus, then whatever's newest.
        chosen = (next((i for i in ids if "sonnet" in i), None)
                  or next((i for i in ids if "opus" in i), None)
                  or (ids[0] if ids else fallback))
        _model_cache["value"] = chosen
        print("AI model auto-selected:", chosen)
        return chosen
    except Exception as e:
        print("model discovery failed (%s) — using fallback %s" % (e, fallback))
        _model_cache["value"] = fallback
        return fallback


# How much deck text we send to the AI. ~25k characters covers virtually every real
# pitch deck (appendix and all) while capping cost on a runaway file. Anything longer
# is flagged to the founder rather than silently dropped.
DECK_CHAR_LIMIT = 25000

def ai_diagnostic(deck_text, icp="", pitch="", images=None):
    if not ANTHROPIC_API_KEY:
        # No key yet -> realistic canned result so the demo flows.
        return {
            "live": False,
            "summary": "Sample result — your real AI key isn't set yet, so this shows the kind of read the engine produces. Add the key and it runs live on the uploaded deck.",
            "findings": [
                {"dim": "Positioning Clarity", "note": "Opening leads with the product, not the customer problem — the 'why now' is missing.", "fixes": [
                    "Open on the customer's painful trigger, then position your product as the answer.",
                    "Add a 'why now' — the market, tech or regulatory shift that makes this the moment.",
                    "Cut feature language from the first slide; lead with the outcome the buyer gets."]},
                {"dim": "ICP Specificity", "note": "Target is described as 'SMEs' with no firmographic or trigger-event qualifier.", "fixes": [
                    "Name the company size and the buyer's exact role (e.g. 'VP Sales at a 50–200-person SaaS').",
                    "State the trigger event that makes them ready to buy right now.",
                    "Name two or three real example companies that fit the profile."]},
                {"dim": "Messaging Hierarchy", "note": "The edge is framed as 'faster and cheaper' — a feature claim, buried among equally-weighted points.", "fixes": [
                    "Lead with the single message that matters most and make everything else support it.",
                    "Replace 'faster and cheaper' with a structural advantage a rival can't copy by trying harder.",
                    "Back the hero claim with one concrete proof point, not three vague ones."]},
                {"dim": "GTM Sequencing", "note": "Acquisition relies on 'word of mouth and partnerships'; there's no first-10-customers plan.", "fixes": [
                    "Spell out exactly how the first ten customers get reached and closed.",
                    "State your motion (self-serve, inside sales, or founder-led) and the CAC you expect.",
                    "Sequence the channels — which one you prove first before layering the next."]},
            ],
        }
    context = ""
    if icp or pitch:
        context = (
            "\n\nThe founder also gave two quick inputs at the free step:\n"
            "WHO IT'S FOR (their stated ICP): " + (icp or "(not given)") + "\n"
            "ONE-LINE PITCH: " + (pitch or "(not given)") + "\n"
            "Use these to judge whether the deck actually delivers on what they claim."
        )
    prompt = (
        "You are a senior TiE Bangalore investment reviewer running a FREE positioning diagnostic. "
        "This free read sits IN FRONT OF a paid expert review, so it must be genuinely useful but must NOT do the "
        "founder's rewriting for them — its job is to make them see WHAT is weak and WHICH DIRECTION to move, while "
        "leaving the actual rework, prioritisation and judgement to the paid expert. "
        "This is AI pattern analysis, NOT a statistical benchmark. Read the founder's pitch deck text below and "
        "identify the most important issue under each of the four positioning dimensions: Positioning Clarity, ICP Specificity, "
        "Messaging Hierarchy, and GTM Sequencing. "
        "For each dimension, the 'note' should be one sharp sentence naming the problem AND why it costs them with investors. "
        "Then give two or three DIRECTIONAL pointers in 'fixes' — short nudges that name the type of fix and why it matters, "
        "phrased as 'rethink / sharpen / pick / separate' guidance. "
        "CRITICAL — do NOT do the work for them: do NOT write the replacement headline or any ready-to-paste copy, "
        "do NOT compute or restate the corrected number, do NOT enumerate the specific certifications, regulations or named "
        "example customers, and do NOT hand over a finished slide. Point the direction only (e.g. 'Pick one beachhead "
        "application and lead the deck with it' — NOT the exact new headline; 'Your cost-reduction claim is mathematically "
        "impossible as stated — restate it as a real from/to figure' — NOT the corrected figure itself). "
        "Watch for these known anti-patterns: " + ANTI_PATTERNS + ". "
        "End 'summary' with one sentence noting that a TiE expert reviewer would help them prioritise which of these fixes "
        "matters most for their specific raise and rework it for the investors they're about to meet. "
        "Respond ONLY as JSON: {\"summary\": str, \"findings\": [{\"dim\": str, \"note\": str, \"fixes\": [str, str, ...]}, ...]}."
        + context + "\n\n"
        + ("The deck is image-based, so its slides are attached as pictures above rather than as text. "
           "Read the slides directly — including any charts, tables and infographics — and diagnose from what you see."
           if images else "DECK TEXT:\n" + deck_text[:DECK_CHAR_LIMIT])
    )
    # Image-only decks: send the slide pictures and let the model read them.
    # Beats OCR — it understands layout, charts and infographics, not just glyphs.
    if images:
        content = [{"type": "image",
                    "source": {"type": "base64", "media_type": im["media_type"], "data": im["b64"]}}
                   for im in images]
        content.append({"type": "text", "text": prompt})
    else:
        content = prompt
    body = json.dumps({
        "model": get_model(),
        # Enough room for a full four-dimension read on a rich, multi-slide deck.
        # Too small and the reply gets cut off mid-JSON and can't be parsed.
        "max_tokens": 4096,
        "messages": [{"role": "user", "content": content}],
    }).encode()
    req = urllib.request.Request(
        "https://api.anthropic.com/v1/messages", data=body, method="POST",
        headers={"x-api-key": ANTHROPIC_API_KEY, "anthropic-version": "2023-06-01", "content-type": "application/json"},
    )
    d = json.loads(urllib.request.urlopen(req, timeout=90).read())
    # Read every text block (newer models can prepend non-text blocks); don't assume
    # the first block is the text — that would crash with a KeyError on those models.
    text = "".join(b.get("text", "") for b in d.get("content", [])
                   if isinstance(b, dict) and b.get("type") == "text")
    start, end = text.find("{"), text.rfind("}")
    try:
        parsed = json.loads(text[start:end + 1])
    except Exception as e:
        # The model's reply wasn't clean JSON (e.g. truncated or wrapped in prose).
        # Fail gracefully with a readable message instead of crashing the request.
        print("AI JSON parse failed (%s); stop_reason=%s" % (e, d.get("stop_reason")))
        return {
            "live": True,
            "partial": True,
            "summary": (
                "The AI read your deck but its detailed reply didn't come back in a "
                "form we could lay out cleanly this time. Please run the diagnostic "
                "again — it usually succeeds on a retry."
            ),
            "findings": [],
        }
    parsed["live"] = True
    return parsed


class Handler(BaseHTTPRequestHandler):
    def _send(self, code, obj):
        payload = json.dumps(obj).encode()
        self.send_response(code)
        self.send_header("Content-Type", "application/json")
        self.send_header("Content-Length", str(len(payload)))
        self.end_headers()
        self.wfile.write(payload)

    def do_GET(self):
        path = "/index.html" if self.path in ("/", "") else self.path.split("?")[0]
        if path == "/api/config":
            # Public front-end config. Only non-secret values belong here. The
            # Google Client ID is designed to be public (the browser needs it).
            return self._send(200, {"googleClientId": GOOGLE_CLIENT_ID,
                                    "membershipUrl": TIE_MEMBERSHIP_URL})
        if path == "/api/support":
            return self._send(200, support_list())
        if path == "/api/review":
            qs = urllib.parse.parse_qs(urllib.parse.urlparse(self.path).query)
            fid = (qs.get("founder_id") or [""])[0]
            return self._send(200, review_state(fid) if fid else {"rounds": [], "latest": None})
        if path == "/api/dataroom":
            qs = urllib.parse.parse_qs(urllib.parse.urlparse(self.path).query)
            fid = (qs.get("founder_id") or [""])[0]
            return self._send(200, dataroom_list(fid) if fid else {"docs": [], "views": []})
        if path == "/api/founder":
            qs = urllib.parse.parse_qs(urllib.parse.urlparse(self.path).query)
            fid = (qs.get("id") or [""])[0]
            rec = founder_get(fid) if fid else None
            return self._send(200, rec or {"error": "not found"})
        if path == "/api/deals":
            # A single party's own deals (their private view).
            qs = urllib.parse.parse_qs(urllib.parse.urlparse(self.path).query)
            fid = (qs.get("founder_id") or [""])[0]
            iem = (qs.get("investor_email") or [""])[0]
            return self._send(200, {"deals": deals_for(fid or None, iem or None)})
        if path == "/api/reviewer/queue":
            qs = urllib.parse.parse_qs(urllib.parse.urlparse(self.path).query)
            em = (qs.get("email") or [""])[0]
            return self._send(200, {"queue": reviewer_queue(em)})
        if path == "/api/impact":
            # PUBLIC, showcase-safe aggregates only (no names, no individual amounts).
            return self._send(200, impact_public())
        fp = os.path.join(HERE, path.lstrip("/"))
        if os.path.isfile(fp) and os.path.abspath(fp).startswith(HERE):
            ctype = "text/html" if fp.endswith(".html") else "application/octet-stream"
            data = open(fp, "rb").read()
            self.send_response(200)
            self.send_header("Content-Type", ctype)
            self.send_header("Content-Length", str(len(data)))
            self.end_headers()
            self.wfile.write(data)
        else:
            self._send(404, {"error": "not found"})

    def do_POST(self):
        length = int(self.headers.get("Content-Length", 0))
        raw = self.rfile.read(length) if length else b"{}"
        try:
            data = json.loads(raw or b"{}")
        except Exception:
            data = {}
        try:
            if self.path == "/api/verify-member":
                _vm = verify_member(data.get("email", ""))
                # A definitive "not a member" (Zoho reached, no active record) is a
                # gate failure worth flagging to admins. "unreachable" is NOT a
                # failure — we just couldn't check — so it does not notify.
                if data.get("email") and (not _vm.get("member")) and (not _vm.get("unreachable")):
                    notify("membership_failed", ["admins"],
                           "Membership check failed at the gate",
                           "A sign-in was blocked: " + data.get("email", "") +
                           " was not found as an active TiE Bangalore member in Zoho.", "high")
                self._send(200, _vm)
            elif self.path == "/api/founder":
                self._send(200, founder_upsert(data))
            elif self.path == "/api/signup":
                out = founder_signup(data)
                self._send(400 if out.get("error") else 200, out)
            elif self.path == "/api/login":
                out = founder_login(data.get("email", ""), data.get("password", ""))
                self._send(401 if out.get("error") else 200, out)
            elif self.path == "/api/auth/google":
                # "Continue with Google": verify the Google credential, find/create
                # the founder, refresh their TiE membership, and return the profile.
                out = founder_google_login(data.get("credential", ""))
                self._send(401 if out.get("error") else 200, out)
            elif self.path == "/api/admin/list":
                # Super-admin only: read the admin allow-list.
                g = require_super_admin(data.get("credential", ""))
                if g.get("error"):
                    self._send(g["status"], {"error": g["error"]})
                else:
                    self._send(200, {"admins": admin_list()})
            elif self.path == "/api/admin/add":
                # Super-admin only: add an approved admin email.
                out = admin_add(data.get("credential", ""), data.get("email", ""))
                if out.get("error"):
                    self._send(out.get("status", 400), {"error": out["error"]})
                else:
                    self._send(200, out)
            elif self.path == "/api/admin/remove":
                # Super-admin only: remove an admin email (super-admin can't be removed).
                out = admin_remove(data.get("credential", ""), data.get("email", ""))
                if out.get("error"):
                    self._send(out.get("status", 400), {"error": out["error"]})
                else:
                    self._send(200, out)
            elif self.path == "/api/admin/notifications":
                # Any admin: view the recorded notifications feed / email outbox.
                g = require_admin(data.get("credential", ""))
                if g.get("error"):
                    self._send(g["status"], {"error": g["error"]})
                else:
                    self._send(200, {"notifications": notifications_recent(),
                                     "mail_enabled": MAIL_ENABLED})
            elif self.path == "/api/review":
                # A reviewer submits a verdict — this is what the founder then sees.
                fid = data.get("founder_id")
                verdict = data.get("verdict", "")
                if not fid or verdict not in ("not_yet", "awarded"):
                    self._send(400, {"error": "founder_id and verdict ('not_yet'|'awarded') are required"})
                else:
                    self._send(200, review_add(fid, verdict, data.get("gaps"),
                                               data.get("reviewer", ""), data.get("note", "")))
            elif self.path == "/api/dataroom":
                # Upload one diligence document against a checklist item.
                fid = data.get("founder_id")
                key = data.get("item_key", "")
                fname = data.get("filename", "")
                b64 = data.get("file_b64", "")
                if not (fid and key and b64):
                    self._send(400, {"error": "founder_id, item_key and file_b64 are required"})
                else:
                    raw = base64.b64decode(b64.split(",")[-1])
                    did = dataroom_add(fid, key, fname, raw)
                    out = dataroom_list(fid)
                    out["doc_id"] = did
                    self._send(200, out)
            elif self.path == "/api/dataroom/view":
                self._send(200, dataroom_view_log(data.get("founder_id"), data.get("doc_id"),
                                                  data.get("item_key", ""), data.get("viewer", "")))
            elif self.path == "/api/diagnostic":
                deck_text = data.get("deck_text", "")
                fileb64 = data.get("file_b64", "")
                fname = data.get("filename", "")
                founder_id = data.get("founder_id")
                had_file = bool(fileb64)
                deck_id = None
                if fileb64:
                    raw = base64.b64decode(fileb64.split(",")[-1])
                    deck_text = extract_deck_text(fname, raw)
                    # Store the actual deck file + a DB row pointing to it.
                    deck_id = deck_add(founder_id, fname, raw)
                # No usable text? The deck is probably image-based (each slide is a
                # picture). Before giving up, try to LOOK at the slides instead.
                deck_images = None
                if had_file and len((deck_text or "").strip()) < 120:
                    deck_images = extract_deck_images(fname, raw)
                    if deck_images:
                        print("image-based deck: reading %d slide picture(s) with vision" % len(deck_images))

                # Only truly unreadable if there's neither text nor a picture to look at.
                if had_file and len((deck_text or "").strip()) < 120 and not deck_images:
                    self._send(200, {
                        "live": False,
                        "unreadable": True,
                        "summary": (
                            "We couldn't read “%s.” There's no text inside it and we "
                            "couldn't pull out readable slide images either. Please export your "
                            "deck as a PDF or PowerPoint and upload it again — or paste your "
                            "pitch below and we'll read that instead."
                        ) % (fname or "your file"),
                        "findings": [],
                    })
                else:
                    result = None
                    # Reuse a prior read for the same founder + identical deck instead of
                    # paying the API again (no double-charging for a re-run).
                    if founder_id and had_file:
                        result = diagnostic_find_reusable(founder_id, fname, len(raw))
                    if result is None:
                        result = ai_diagnostic(deck_text, data.get("icp", ""), data.get("pitch", ""),
                                               images=deck_images)
                        if deck_images:
                            # Be transparent that this read came from looking at the slides.
                            result["read_from_images"] = True
                            result["summary"] = (
                                "Note: this deck has no text inside it, so we read the slides as images. "
                            ) + (result.get("summary") or "")
                        # Persist the read against this founder + deck so every screen and a
                        # later visit can pull it back up (survives refresh and restart).
                        if founder_id:
                            diagnostic_save(founder_id, deck_id, result)
                    # Long-deck honesty flag: if the deck ran past what we read in full,
                    # tell the founder rather than silently analysing only the first part.
                    if had_file and len(deck_text or "") > DECK_CHAR_LIMIT:
                        result["long_deck"] = True
                        result["summary"] = (
                            "Heads-up: this deck is longer than we read in full — only the first "
                            "~25,000 characters were analysed, so the tail end wasn't seen. Consider "
                            "trimming or splitting it for a complete read. "
                        ) + (result.get("summary") or "")
                    result["deck_id"] = deck_id
                    self._send(200, result)
            elif self.path == "/api/reviewer/signup":
                out = reviewer_signup(data)
                self._send(400 if out.get("error") else 200, out)
            elif self.path == "/api/reviewer/login":
                out = reviewer_login(data.get("email", ""), data.get("password", ""))
                self._send(401 if out.get("error") else 200, out)
            elif self.path == "/api/reviewer/auth/google":
                out = reviewer_google_login(data.get("credential", ""))
                self._send(401 if out.get("error") else 200, out)
            elif self.path == "/api/reviewer/submit":
                out = reviewer_submit(data.get("reviewer_email", ""), data.get("founder_id"),
                                      data.get("verdict", ""), data.get("gaps"), data.get("note", ""))
                self._send(out.get("status", 400) if out.get("error") else 200, out)
            elif self.path == "/api/admin/pending-verdicts":
                out = pending_verdicts(data.get("credential", ""))
                if out.get("error"):
                    self._send(out.get("status", 403), {"error": out["error"]})
                else:
                    self._send(200, out)
            elif self.path == "/api/review/countersign":
                out = review_countersign(data.get("credential", ""), data.get("review_id"),
                                         data.get("decision", ""), data.get("note", ""))
                if out.get("error"):
                    self._send(out.get("status", 400), {"error": out["error"]})
                else:
                    self._send(200, out)
            elif self.path == "/api/deal/create":
                out = deal_create(data)
                self._send(out.get("status", 400) if out.get("error") else 200, out)
            elif self.path == "/api/deal/advance":
                out = deal_advance(data)
                self._send(out.get("status", 400) if out.get("error") else 200, out)
            elif self.path == "/api/deal/verify":
                out = deal_verify(data.get("credential", ""), data.get("deal_id"),
                                  data.get("verified", True))
                if out.get("error"):
                    self._send(out.get("status", 400), {"error": out["error"]})
                else:
                    self._send(200, out)
            elif self.path == "/api/admin/deals":
                out = deals_admin(data.get("credential", ""))
                if out.get("error"):
                    self._send(out.get("status", 403), {"error": out["error"]})
                else:
                    self._send(200, out)
            elif self.path == "/api/support":
                self._send(200, support_add(data))
            elif self.path == "/api/payment/record":
                out = payment_record(data)
                self._send(out.get("status", 400) if out.get("error") else 200, out)
            elif self.path == "/api/admin/review-queue":
                g = require_admin(data.get("credential", ""))
                if g.get("error"):
                    self._send(g["status"], {"error": g["error"]})
                else:
                    self._send(200, {"queue": review_queue()})
            elif self.path == "/api/review/assign":
                out = review_assign(data.get("credential", ""), data.get("founder_id"),
                                    data.get("reviewer_email", ""), data.get("reviewer_name", ""))
                if out.get("error"):
                    self._send(out.get("status", 400), {"error": out["error"]})
                else:
                    self._send(200, out)
            elif self.path == "/api/investor/apply":
                out = investor_apply(data)
                self._send(out.get("status", 400) if out.get("error") else 200, out)
            elif self.path == "/api/admin/investors":
                out = investor_list(data.get("credential", ""))
                if out.get("error"):
                    self._send(out.get("status", 403), {"error": out["error"]})
                else:
                    self._send(200, out)
            elif self.path == "/api/investor/decision":
                out = investor_decision(data.get("credential", ""), data.get("investor_id"),
                                        data.get("decision", ""), data.get("reason", ""))
                if out.get("error"):
                    self._send(out.get("status", 400), {"error": out["error"]})
                else:
                    self._send(200, out)
            elif self.path == "/api/membership/handoff":
                self._send(200, membership_handoff(data))
            elif self.path == "/api/membership/recheck":
                out = membership_recheck(data)
                self._send(out.get("status", 400) if out.get("error") else 200, out)
            elif self.path == "/api/investor/invite":
                out = investor_invite(data)
                self._send(out.get("status", 400) if out.get("error") else 200, out)
            elif self.path == "/api/meeting/request":
                out = meeting_request(data)
                self._send(out.get("status", 400) if out.get("error") else 200, out)
            elif self.path == "/api/admin/meetings":
                out = meeting_list(data.get("credential", ""))
                if out.get("error"):
                    self._send(out.get("status", 403), {"error": out["error"]})
                else:
                    self._send(200, out)
            elif self.path == "/api/meeting/intro-sent":
                out = meeting_intro_sent(data.get("credential", ""), data.get("meeting_id"))
                if out.get("error"):
                    self._send(out.get("status", 400), {"error": out["error"]})
                else:
                    self._send(200, out)
            else:
                self._send(404, {"error": "unknown endpoint"})
        except Exception as e:
            self._send(500, {"error": str(e)})

    def log_message(self, *args):
        pass  # quiet console


if __name__ == "__main__":
    port = 8000
    db_init()
    print("VentureReady demo running.  Open this in your browser:  http://localhost:%d" % port)
    print("Zoho member check: %s | AI key: %s" % (
        "ready" if ZOHO_CLIENT_ID else "MISSING",
        "set (live AI)" if ANTHROPIC_API_KEY else "not set (sample result)"))
    print("Press Ctrl+C here to stop the app.")
    ThreadingHTTPServer(("127.0.0.1", port), Handler).serve_forever()
